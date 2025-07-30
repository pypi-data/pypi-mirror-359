# %% Build file search index

## Load dependencies

import re
import os
import csv
import sys
import glob
import email
import PyPDF2
import shutil
import os.path
import hashlib
import logging
import zipfile
import getpass
import openpyxl
import textwrap
import numpy as np
import polars as pl
import pandas as pd
import subprocess
from pathlib import Path
from ebooklib import epub
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
from extract_msg import Message
from dateutil.parser import parse
from email.utils import parseaddr
from datetime import datetime, timezone
import concurrent.futures
from analytics_tasks.utils.functions import timer_start, timer_end

import os
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor
import multiprocessing as mp


## Assign folder or file names
folder_dt = datetime.now(timezone.utc).strftime("%Y%m%d")
file_dt = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")


# %% build_process

## Declare a global variable at the module level
time_machine = None


## 0. lib_refs_fs
def lib_refs_fs(at_dir):
    """Assigns working libraries inside file_search dir."""
    timer_start()
    _logs_dir = at_dir / "file_search/output/log"
    _time_machine_path = at_dir / "file_search/input/time_machine"
    _reports_dir = at_dir / "file_search/output/report"

    Path(_logs_dir).mkdir(parents=True, exist_ok=True)
    Path(_time_machine_path).mkdir(parents=True, exist_ok=True)
    Path(_reports_dir).mkdir(parents=True, exist_ok=True)

    # Fixed folder to keep search index
    _fs_index_dir = at_dir / "file_search/input"

    # Assign the path to the global variable
    global time_machine
    time_machine = _time_machine_path

    print("Assigned file_search directories.")

    # return _fs_index_dir, _logs_dir, _time_machine_path, _reports_dir
    return _fs_index_dir, _logs_dir, _time_machine_path, _reports_dir


# %% functions


def wt(table, folder_level, file_group):
    """calculate weight of folder(s) by level"""

    global levelx

    try:
        del levelx
    except Exception:
        pass

    levelx = table[["uf_id", folder_level, "size_mb"]].drop_duplicates()
    levelx = (
        levelx.groupby([folder_level])
        .agg({"uf_id": "nunique", "size_mb": "sum"})
        .rename(columns={"uf_id": "#of files"})
        .sort_values(by="size_mb", ascending=False)
        .reset_index()
    )

    sl = levelx[r"#of files"].sum()
    if sl != 0:
        _s1 = f"NOTE: Importing files = {str(sl)}"
        print("\n")
        print("-" * len(_s1))
        print(_s1)
        print("PROCESSING:", str(file_group))
        print(r"SIZE: " + str(round(levelx.size_mb.sum(), 3)) + r" MB")
        print("-" * len(_s1))
        print(levelx)
        print("\n")
    else:
        print("WARNING: No new files to scan in", file_group)


def fu(table, time_field, folder_level):
    """calculate folder usage by folder level"""

    global folder_usage
    global folder_usagex

    try:
        del folder_usage
        del folder_usagex
    except Exception:
        print("\t")

    folder_usage = (
        table.groupby([time_field, "ext", folder_level])
        .agg({"unc": "nunique"})
        .reset_index()
        .sort_values(by=[time_field, "unc"], ascending=[False, False])
    )
    folder_usagex = (
        folder_usage.groupby([time_field, "ext"])
        .agg({"unc": "sum"})
        .reset_index()
        .sort_values(by=[time_field, "unc"], ascending=[False, False])
    )
    folder_usagex = (
        folder_usagex.pivot(time_field, "ext", "unc")
        .sort_values(by=[time_field])
        .fillna("0")
    )

    for i in folder_usagex.columns:
        folder_usagex.loc[:, i] = (
            pd.to_numeric(folder_usagex.loc[:, i], errors="coerce")
            .round()
            .astype("int")
        )
    folder_usagex = folder_usagex.replace(0, ".").sort_values(
        by=[time_field], ascending=False
    )
    folder_usagex = folder_usagex.reset_index()

    print(
        r"\nNOTE: UNC count by folder level",
        folder_level,
        r" and",
        time_field,
        r"\n",
        folder_usage,
    )
    print(r"\nNOTE: UNC count by", time_field, r" \n", folder_usagex)


def scan_folder_searchx(location_to_scan=time_machine, *, ext=""):
    """scan a path for file types"""

    scan = []
    if ext == "":
        for i in glob.iglob(f"{location_to_scan}\\*", recursive=True):
            scan.append(i)
    else:
        for i in glob.iglob(
            f"{location_to_scan}\\**\\*{ext}".format(ext), recursive=True
        ):
            scan.append(i)
    if len(scan) > 0:
        scan = pd.DataFrame(scan).rename(columns={0: "unc"})
        scan["unc"] = scan["unc"].str.replace("\\", "/")
        scan["uf_id"] = (
            scan["unc"]
            .apply(lambda row: os.path.splitext(os.path.basename(row))[0])
            .astype("int64")
        )
    else:
        scan = pd.DataFrame(
            columns=[
                "owner",
                "unc",
                "creationtime",
                "creationtimeutc",
                "lastaccesstime",
                "lastaccesstimeutc",
                "lastwritetime",
                "lastwritetimeutc",
                "filename",
                "ext",
                "size_mb",
                "uf_id",
            ]
        )

    return scan


def hash_file_info(row):
    """create hash id for file unc"""
    _unc = row["unc"]
    _length = row["length"]
    _lastwritetimeutc = row["lastwritetimeutc"]

    # Get file size and last modification time
    file_size = str(_length)
    modification_time = str(_lastwritetimeutc)

    # Concatenate file path, size, and last modification time
    combined_info = f"{_unc}{file_size}{modification_time}"

    # Hash the combined information
    hash_object = hashlib.sha256()
    hash_object.update(combined_info.encode("utf-8"))

    return int(hash_object.hexdigest(), 16) % (
        10**8
    )  # Module to keep it a small number


def analyze_imoprt_load(scan, scan_ext, file_group):
    """understand the bulk of files to be imported"""
    df = scan[scan["ext"].isin(scan_ext[file_group])]
    wt(df, "ext", file_group)
    df = df[["unc", "uf_id"]].reset_index(drop=True)
    return df


def read_text(unc):
    """read text files"""
    try:
        with open(unc, "r", encoding="utf-8") as file:
            lines = file.readlines()
    except Exception as e:
        print(unc)
        print(f"Error while reading the file: {e}")
        lines = []

    df = pd.DataFrame({"text": lines})
    return df


def read_docx(file_path):
    """read .docx files"""
    doc = Document(file_path)
    paragraphs = []
    for paragraph in doc.paragraphs:
        paragraphs.append(paragraph.text)

    df = pd.DataFrame({"text": paragraphs})
    return df


# %% read_pptx


def read_pptx_old(file_path):
    """read .pptx files"""
    data = []
    presentation = Presentation(file_path)
    for slide_num, slide in enumerate(presentation.slides, start=1):
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

        data.append({"page": slide_num, "text": text})
    df = pd.DataFrame(data)
    return df


def read_pptx(file_path):
    """read .pptx files with position information and exploded rows for each text element"""
    all_text_elements = []
    presentation = Presentation(file_path)

    for slide_num, slide in enumerate(presentation.slides, start=1):
        texts_with_positions = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Get position information
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                texts_with_positions.append(
                    {
                        "text": shape.text,
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                    }
                )

        # Sort texts by vertical position (top to bottom)
        texts_with_positions.sort(key=lambda x: x["top"])

        # Create individual rows for each text element
        for i, text_element in enumerate(texts_with_positions):
            # First n elements are flagged as 1, the rest as 0
            is_top_comment = 1 if i < 3 else 0

            all_text_elements.append(
                {
                    "page": slide_num,
                    "text": text_element["text"],
                    "left": text_element["left"],
                    "top": text_element["top"],
                    "width": text_element["width"],
                    "height": text_element["height"],
                    "comments_top": is_top_comment,
                    "element_position": i,  # Add the position in the slide for reference
                }
            )

    df = pd.DataFrame(all_text_elements)
    df = df[["page", "text", "comments_top"]]

    return df


# %% read_msg


def read_msg(file_path):
    """read .msg files"""
    msg = Message(file_path)

    subject = msg.subject
    sender = msg.sender
    recipients = msg.to
    sent_date = pd.to_datetime(msg.date, utc=None)
    body = msg.body

    attachmentsx = [attachment.longFilename for attachment in msg.attachments]
    attachments = "|".join(map(str, attachmentsx))

    df = pd.DataFrame(
        {
            "subject": [subject],
            "sender": [sender],
            "recipients": [recipients],
            "sentdate": [sent_date],
            "text": [body],
            "attachments": [attachments],
        }
    )
    return df


# %% read_eml


def read_eml(file_path):
    """Read .eml files"""
    with open(file_path, "rb") as file:
        msg = email.message_from_binary_file(file)

    subject = msg["subject"]
    sender = parseaddr(msg["from"])[1]  # Extract email address
    recipients = msg["to"]
    sent_date = pd.to_datetime(msg["date"], utc=True)

    # Get the body of the email
    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body = part.get_payload(decode=True).decode()
                break
    else:
        body = msg.get_payload(decode=True).decode()

    # Get attachments
    attachmentsx = []
    for part in msg.walk():
        if part.get_content_maintype() == "multipart":
            continue
        if part.get("Content-Disposition") is None:
            continue
        filename = part.get_filename()
        if filename:
            attachmentsx.append(filename)

    attachments = "|".join(map(str, attachmentsx))

    df = pd.DataFrame(
        {
            "subject": [subject],
            "sender": [sender],
            "recipients": [recipients],
            "sentdate": [sent_date],
            "text": [body],
            "attachments": [attachments],
        }
    )

    return df


# %% read_epub


def read_epub(file_path):
    """read .epub files"""
    book = epub.read_epub(file_path)

    title = (
        book.get_metadata("DC", "title")[0][0]
        if book.get_metadata("DC", "title")
        else None
    )
    author = (
        book.get_metadata("DC", "creator")[0][0]
        if book.get_metadata("DC", "creator")
        else None
    )
    language = (
        book.get_metadata("DC", "language")[0][0]
        if book.get_metadata("DC", "language")
        else None
    )
    identifier = (
        book.get_metadata("DC", "identifier")[0][0]
        if book.get_metadata("DC", "identifier")
        else None
    )

    content = ""
    for item in book.get_items_of_type(9):  # type 9 corresponds to text content
        content += item.content.decode("utf-8", "ignore")

    soup = BeautifulSoup(content, "html.parser")
    cleaned_text_lines = list(
        filter(None, map(str.strip, soup.get_text().splitlines()))
    )

    df = pd.DataFrame({"text": [cleaned_text_lines]})

    df = df.explode("text").reset_index(drop=True)

    return df


# %% read_pdf


def read_pdf(file_path):
    """read .pdf files"""
    data = []

    with open(file_path, "rb") as file:
        pdf_reader = PyPDF2.PdfReader(file)

        for page_num, page in enumerate(pdf_reader.pages, start=1):
            text = page.extract_text()

            lines = text.split("\n")

            for line_num, line in enumerate(lines, start=1):
                data.append({"page": page_num, "row": line_num, "text": line})

    df = pd.DataFrame(data)
    return df


# %% read_xlsx


def read_xlsx_formula(unc):
    """read formula in .xlsx files"""
    row_formula_df = pd.DataFrame()

    workbook = openpyxl.load_workbook(unc, data_only=False)

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        df = pd.DataFrame(sheet.values)

        rows_formulas = df.map(lambda x: isinstance(x, str) and x.startswith("="))
        indices = [
            (index, column)
            for index, series in rows_formulas.iterrows()
            for column, value in series.items()
            if value
        ]
        row_formula_dfx = pd.DataFrame(indices, columns=["row", "column"])
        if len(row_formula_dfx) != 0:
            try:
                row_formula_dfx["formula_full"] = row_formula_dfx.apply(
                    lambda x: df.at[x["row"], x["column"]], axis=1
                )

                row_formula_dfx["function"] = row_formula_dfx[
                    "formula_full"
                ].str.extract(r"=([A-Za-z_]\w*)", expand=False)
                row_formula_dfx["sheet"] = sheet_name
                row_formula_dfx["column"] = row_formula_dfx["column"] + 1
                row_formula_df = pd.concat([row_formula_df, row_formula_dfx])
                workbook.close()

            except ValueError:
                continue

    if len(row_formula_df) == 0:
        row_formula_df = pd.DataFrame(columns=["row", "column", "sheet"])
    return row_formula_df


def read_xlsx_text(unc):
    """read text in .xlsx files"""
    workbook = openpyxl.load_workbook(unc, data_only=True)
    all_data = []

    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        max_row = sheet.max_row
        max_col = sheet.max_column
        data = []
        for row in range(1, max_row + 1):
            for col in range(1, max_col + 1):
                cell_value = sheet.cell(row=row, column=col).value

                if isinstance(cell_value, str) and not cell_value.isdigit():
                    data.append([row, col, cell_value, sheet_name])
        all_data.extend(data)

    df = pd.DataFrame(all_data, columns=["row", "column", "text", "sheet"])
    workbook.close()

    if len(df) == 0:
        df = pd.DataFrame(columns=["row", "column", "text", "sheet"])

    return df


# %% load_ifp


# def load_ifp(scan, scan_ext, file_group, _function):
#     """function to load files belonging to a particular group"""

#     select_file_group = analyze_imoprt_load(scan, scan_ext, file_group)
#     r = 0
#     range_max = len(select_file_group)
#     for unc, uf_id in select_file_group.itertuples(index=False):
#         print("READING:", unc)
#         r += 1
#         try:
#             dfx = _function(unc)
#             if dfx.isna().all().all():
#                 print("WARNING: file does not have any content")
#                 dfx = pd.DataFrame({"text": "error reading file"}, index=[0])
#                 dfx.to_pickle(
#                     str(time_machine).replace("\\", "/") + "/" + str(uf_id) + ".pickle"
#                 )
#             else:
#                 dfx.to_pickle(
#                     str(time_machine).replace("\\", "/") + "/" + str(uf_id) + ".pickle"
#                 )
#                 print("IMPORTED:", f"{r:09}", "of", f"{range_max:09}", unc)
#         except Exception:
#             dfx = pd.DataFrame({"text": "error reading file"}, index=[0])
#             dfx.to_pickle(
#                 str(time_machine).replace("\\", "/") + "/" + str(uf_id) + ".pickle"
#             )


def load_ifp(scan, scan_ext, file_group, _function):
    """function to load files belonging to a particular group"""

    select_file_group = analyze_imoprt_load(scan, scan_ext, file_group)
    r = 0
    range_max = len(select_file_group)

    for unc, uf_id in select_file_group.itertuples(index=False):
        print("READING:", unc)
        r += 1

        try:
            # Use ThreadPoolExecutor with timeout
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                # Submit the function
                future = executor.submit(_function, unc)

                try:
                    # Wait for result with 4 minute timeout
                    dfx = future.result(timeout=240)

                    if dfx.isna().all().all():
                        print("WARNING: file does not have any content")
                        dfx = pd.DataFrame({"text": "error reading file"}, index=[0])
                        dfx.to_pickle(
                            str(time_machine).replace("\\", "/")
                            + "/"
                            + str(uf_id)
                            + ".pickle"
                        )
                    else:
                        dfx.to_pickle(
                            str(time_machine).replace("\\", "/")
                            + "/"
                            + str(uf_id)
                            + ".pickle"
                        )
                        print("IMPORTED:", f"{r:09}", "of", f"{range_max:09}", unc)

                except concurrent.futures.TimeoutError:
                    print(
                        f"TIMEOUT: Function took too long for file {unc}, skipping..."
                    )
                    # Cancel the future (though this might not stop the actual execution)
                    future.cancel()

                    dfx = pd.DataFrame(
                        {"text": "error reading file - timeout"}, index=[0]
                    )
                    dfx.to_pickle(
                        str(time_machine).replace("\\", "/")
                        + "/"
                        + str(uf_id)
                        + ".pickle"
                    )

        except Exception:
            dfx = pd.DataFrame({"text": "error reading file"}, index=[0])
            dfx.to_pickle(
                str(time_machine).replace("\\", "/") + "/" + str(uf_id) + ".pickle"
            )


# def load_ifp_xlsx(scan, scan_ext, file_group, _function1, _function2):
#     """function to load files belonging to a particular group"""

#     select_file_group = analyze_imoprt_load(scan, scan_ext, file_group)
#     r = 0
#     range_max = len(select_file_group)
#     for unc, uf_id in select_file_group.itertuples(index=False):
#         r += 1
#         try:
#             df1 = _function1(unc)
#             df2 = _function2(unc)
#             if len(df2) == 0:
#                 dfx = df1.copy()
#             else:
#                 dfx = pd.merge(df1, df2, how="outer", on=["row", "column", "sheet"])

#             if dfx.isna().all().all():
#                 print("WARNING: file does not have any content")
#                 dfx = pd.DataFrame({"text": "error reading file"}, index=[0])
#                 dfx.to_pickle(
#                     str(time_machine).replace("\\", "/") + "/" + str(uf_id) + ".pickle"
#                 )
#             else:
#                 dfx.to_pickle(
#                     str(time_machine).replace("\\", "/") + "/" + str(uf_id) + ".pickle"
#                 )
#                 print("\nIMPORTED:", f"{r:09}", "of", f"{range_max:09}", unc)
#         except Exception:
#             dfx = pd.DataFrame({"text": "error reading file"}, index=[0])
#             dfx.to_pickle(
#                 str(time_machine).replace("\\", "/") + "/" + str(uf_id) + ".pickle"
#             )
#             print("ERROR: function load_ifp_xlsx... cannot read file", unc)


def load_ifp_xlsx(scan, scan_ext, file_group, _function1, _function2):
    """function to load files belonging to a particular group"""

    select_file_group = analyze_imoprt_load(scan, scan_ext, file_group)
    r = 0
    range_max = len(select_file_group)

    for unc, uf_id in select_file_group.itertuples(index=False):
        print("READING:", unc)
        r += 1
        try:
            # Use ThreadPoolExecutor with timeout for each function separately
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                # Submit first function
                future1 = executor.submit(_function1, unc)

                try:
                    # Wait for first function result with 4 minute timeout
                    df1 = future1.result(timeout=240)
                except concurrent.futures.TimeoutError:
                    print(
                        f"TIMEOUT: Function1 took too long for file {unc}, skipping..."
                    )
                    future1.cancel()
                    dfx = pd.DataFrame(
                        {"text": "error reading file - timeout"}, index=[0]
                    )
                    dfx.to_pickle(
                        str(time_machine).replace("\\", "/")
                        + "/"
                        + str(uf_id)
                        + ".pickle"
                    )
                    continue

            # If first function succeeded, run second function
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                # Submit second function
                future2 = executor.submit(_function2, unc)

                try:
                    # Wait for second function result with 4 minute timeout
                    df2 = future2.result(timeout=240)
                except concurrent.futures.TimeoutError:
                    print(
                        f"TIMEOUT: Function2 took too long for file {unc}, skipping..."
                    )
                    future2.cancel()
                    dfx = pd.DataFrame(
                        {"text": "error reading file - timeout"}, index=[0]
                    )
                    dfx.to_pickle(
                        str(time_machine).replace("\\", "/")
                        + "/"
                        + str(uf_id)
                        + ".pickle"
                    )
                    continue

            # Both functions completed successfully
            if len(df2) == 0:
                dfx = df1.copy()
            else:
                dfx = pd.merge(df1, df2, how="outer", on=["row", "column", "sheet"])

            if dfx.isna().all().all():
                print("WARNING: file does not have any content")
                dfx = pd.DataFrame({"text": "error reading file"}, index=[0])
                dfx.to_pickle(
                    str(time_machine).replace("\\", "/") + "/" + str(uf_id) + ".pickle"
                )
            else:
                dfx.to_pickle(
                    str(time_machine).replace("\\", "/") + "/" + str(uf_id) + ".pickle"
                )
                print("IMPORTED:", f"{r:09}", "of", f"{range_max:09}", unc)

        except Exception as e:
            print("ERROR: function load_ifp_xlsx... cannot read file", unc)
            print(f"Error details: {e}")
            dfx = pd.DataFrame({"text": "error reading file"}, index=[0])
            dfx.to_pickle(
                str(time_machine).replace("\\", "/") + "/" + str(uf_id) + ".pickle"
            )


def clean_import_load(df):
    """Remove null and unwanted values from text."""
    try:
        df["text"] = df["text"].str.replace("\n", "")
        df = df.replace("nan", np.nan)
        df = df[df["text"] != ""]
        df = df[df["text"] != " "]
        df = df[df["text"] != "."]
        df = df[~df["text"].isnull()]
        df = df.reset_index(drop=True)
    except Exception:
        df = pd.DataFrame({"text": "", "unc": ""}, index=[0])
    return df


def split_text(df, max_length=256):
    new_rows = []

    for _, row in df.iterrows():
        text = row["text"]

        if not isinstance(text, str):
            text = str(text) if pd.notna(text) else ""

        # Use textwrap to split the text while respecting word boundaries
        wrapped_text = textwrap.fill(text, width=max_length)

        # Split the wrapped text into lines
        lines = wrapped_text.split("\n")

        for line in lines:
            new_row = {col: row[col] for col in df.columns}
            new_row["text"] = line
            new_rows.append(new_row)

    new_df = pd.DataFrame(new_rows, columns=df.columns)
    return new_df


def remove_emojis(string):
    emoji_pattern = re.compile(
        "["
        "\U0001f600-\U0001f64f"  # emoticons
        "\U0001f300-\U0001f5ff"  # symbols & pictographs
        "\U0001f680-\U0001f6ff"  # transport & map symbols
        "\U0001f1e0-\U0001f1ff"  # flags (iOS)
        "\U00002702-\U000027b0"
        "\U000024c2-\U0001f251"
        "]+",
        flags=re.UNICODE,
    )

    return emoji_pattern.sub(r"", string)


def unzip_file(zip_path, output_folder):
    with zipfile.ZipFile(zip_path, "r") as zip_ref:
        zip_ref.extractall(output_folder)


def delete_file(file_path):
    try:
        os.remove(file_path)
        print(f"NOTE: File '{file_path}' deleted successfully.")
    except OSError as e:
        print(f"ERROR: {file_path} : {e.strerror}")


def zip_folder(folder_path, output_path):
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zipf:
        for root, _, files in os.walk(folder_path):
            for file in files:
                abs_path = os.path.join(root, file)
                rel_path = os.path.relpath(abs_path, folder_path)
                zipf.write(abs_path, rel_path)


def delete_folder(folder_path):
    try:
        shutil.rmtree(folder_path)
        print(f"NOTE: Folder '{folder_path}' and its contents deleted successfully.")
    except OSError as e:
        print(f"ERROR: {folder_path} : {e.strerror}")


def log_start(folder_location):
    """Start logging process"""

    global file_handler, __log_name

    # Check if the logger is already set up
    if "file_handler" in globals() and file_handler is not None:
        print(
            "WARNING : Logging is already in progress. Call log_end() before starting a new log."
        )
        return

    class LogPrints:
        def __init__(self, logger, level=logging.INFO):
            self.logger = logger
            self.level = level
            self.linebuf = ""

        def write(self, buf):
            self.linebuf += buf
            lines = self.linebuf.split("\n")
            for line in lines[:-1]:
                self.logger.log(self.level, line)
            self.linebuf = lines[-1]

        def flush(self):
            pass

    # Define file_dt globally for demonstration purposes
    file_dt = datetime.now().strftime("%Y%m%d_%H%M%S")

    _log_name = "log_" + file_dt + ".log"
    __log_name = str(folder_location) + "\\" + _log_name

    # Create a logger only if it doesn't exist
    if "logger" not in globals():
        logger = logging.getLogger(str(getpass.getuser()))
        logger.setLevel(logging.DEBUG)

        # Create a console handler and set the level to DEBUG
        console_handler = logging.StreamHandler()
        console_handler.setLevel(logging.DEBUG)

        # Removes previously set handlers
        logger.handlers = []

        # Create a file handler and add it to the logger
        file_handler = logging.FileHandler(__log_name, "w+", encoding="utf-8")

        # Formatting
        formatter = logging.Formatter(
            "%(asctime)s | [%(levelname)s] | %(name)s | %(message)s"
        )
        file_handler.setFormatter(formatter)

        formatter_console = logging.Formatter("%(message)s")
        console_handler.setFormatter(formatter_console)

        # Add the handlers to the logger
        logger.addHandler(file_handler)
        logger.addHandler(console_handler)

        # Redirect prints to the log
        sys.stdout = LogPrints(logger, level=logging.INFO)
        sys.stderr = LogPrints(logger, level=logging.ERROR)

        print("NOTE: Logging started...", __log_name)
    else:
        print("NOTE: Logging is already in progress.")


# %% log_end


def log_end():
    """end logging process"""
    global file_handler

    if file_handler is not None:
        print("NOTE: Logging ended...", __log_name)

        # Close the file handler (this will also flush the log entries to the file)
        file_handler.flush()
        file_handler.close()

        # Reset stdout and stderr to their original values
        sys.stdout = sys.__stdout__
        sys.stderr = sys.__stderr__

        # Set file_handler to None to indicate that logging is not in progress
        file_handler = None
    else:
        print("\nNOTE: Logging is not in progress. No action taken.")


# %% scan_dir_powershell


def scan_dir_powershell(fs_index_dir, scan_dirs):
    # Scan directories with powershell
    powershell_file = fs_index_dir / "scan.ps1"
    f = open(fs_index_dir / "scan.ps1", "w")
    for i in range(0, len(scan_dirs)):
        print(f"NOTE: writing powershell script to: {powershell_file}", i)
        main_string = (
            "\nget-childitem "
            + "'"
            + str(scan_dirs[i])
            + "'"
            + " -recurse | where {!$_.PSIsContainer} | select-object @{N='Owner';E={$_.GetAccessControl().Owner}}, FullName, PSChildName, CreationTime, CreationTimeUtc, LastAccessTime, LastAccessTimeUtc, LastWriteTime, LastWriteTimeUtc, Length, Extension | export-csv -notypeinformation -delimiter "
            + "'"
            + "|"
            + "'"
            + " -path "
            + '"'
            + str(fs_index_dir)
            + r"\scan.csv"
            + '" -Encoding UTF8'
        )
        if i == 0:
            f.write("cd " + "'" + str(fs_index_dir) + "'" + "\n")
            f.write(main_string)
        else:
            f.write(main_string + " -append")
    f.close()

    # Run powershell scan
    try:
        print("NOTE: Running powershell script...")
        subprocess.check_output(
            "powershell -Executionpolicy ByPass -file {}".format("scan.ps1")
        )
    except Exception:
        print(
            "ERROR: Powershell script did not run. "
            f"Please run {powershell_file} manually."
        )


# %% time machine


def scan_time_machine(_time_machine_path):
    """Checks past index availability."""
    try:
        if (os.path.isfile(str(_time_machine_path) + ".zip") == 0) & (
            os.path.isdir(_time_machine_path)
        ):
            print("\nNOTE: Directory time_machine exists.")
        elif os.path.isfile(str(_time_machine_path) + ".zip"):
            print("\nREPORT: Loading previous state...")
            # Unzip time_machine.zip
            unzip_file(str(_time_machine_path) + ".zip", _time_machine_path)

            # Delete time_machine.zip
            delete_file(str(_time_machine_path) + ".zip")
    except Exception:
        print(
            "NOTE: Time_machine not present, ignore if building index for the first time."
        )
        pass


# %% Scan clean


def scan_clean(fs_index_dir):
    # import
    scan0 = pd.read_csv(str(fs_index_dir) + "\\scan.csv", sep="|")

    # rename
    scan0.columns = map(str.lower, scan0.columns)
    scan0 = scan0.rename(columns={"fullname": "unc"})

    ## derive: filename, ext, size_mb

    # filename
    scan0["filename"] = scan0["unc"].apply(lambda row: Path(row).name)

    # extension
    scan0["ext"] = scan0["unc"].apply(
        lambda row: os.path.splitext(os.path.basename(row))[1]
    )

    # size: MB
    scan0["size_mb"] = scan0.apply(lambda row: row.length / (1024 * 1024), axis=1)

    ## derive unique file id
    scan0["uf_id"] = scan0.apply(hash_file_info, axis=1)

    # format datetime
    # List of date formats to try
    date_formats = [
        "%m-%d-%Y",
        "%m/%d/%Y %H:%M:%S %p",
        "%Y-%m-%dT%H:%M:%S",
        "%Y/%m/%dT%H:%M:%S",
    ]

    # Function to attempt datetime conversion with multiple formats
    def convert_datetime(date_str):
        for fmt in date_formats:
            try:
                return pd.to_datetime(date_str, format=fmt)
            except ValueError:
                pass
        # Fallback using dateutil parser
        try:
            return parse(date_str)
        except (ValueError, TypeError):
            return pd.NaT

    # Apply the conversion function to each datetime column
    datetime_columns = [
        "creationtime",
        "creationtimeutc",
        "lastaccesstime",
        "lastaccesstimeutc",
        "lastwritetime",
        "lastwritetimeutc",
    ]

    for col in datetime_columns:
        scan0[col] = scan0[col].apply(convert_datetime)

    # drop columns
    scan0 = scan0.drop(["length", "pschildname", "extension"], axis=1)

    # sort
    scan0 = scan0.sort_values(by="lastwritetimeutc", ascending=False)
    scan0 = scan0.reset_index(drop=True)

    return scan0


# %% Previous scan


def scan_history(scan0, _fs_index_dir):
    """If no previous scan, starts everything anew."""
    os.chdir(_fs_index_dir)
    try:
        scan_old = pd.read_parquet("scan0.parquet")
        searchx_final_old = pd.read_parquet("searchx_final.parquet")

        # filter files different from previous scan
        common_files = scan0[scan0["uf_id"].isin(searchx_final_old["uf_id"].tolist())]

        # same as last time
        scan_old = scan_old[scan_old["unc"].isin(common_files.unc.unique())]
        searchx_final_old = searchx_final_old[
            searchx_final_old["uf_id"].isin(common_files["uf_id"].unique())
        ]

        # files different from last time
        scan = scan0[~scan0["unc"].isin(common_files.unc.unique())]
        scan = scan.reset_index(drop=True)

        print("\nREPORT: Total new files to be scanned is", len(scan.unc.unique()))

    except Exception:
        print("REPORT: No previous scan backup found.")
        print(
            "WARNING: Run may take a long time to complete depending on the disk speed and size."
        )
        scan = scan0.copy()
        scan_old = pd.DataFrame()
        searchx_final_old = pd.DataFrame({"uf_id": "", "text": ""}, index=[0])

    return scan, scan_old, searchx_final_old


# %% Exceptions


## Paths to exclude
def exceptions(scan, paths_to_exclude):
    # exclude paths
    print("REPORT: Files to scan before applying exceptions", len(scan))

    exception_path = pd.DataFrame()
    print("NOTE: path_to_exclude...")
    for i in paths_to_exclude:
        print(f"\t {str(i)}")
        exception_pathx = scan[
            scan["unc"]
            .str.replace("\\", "/")
            .str.contains(str(i).replace("\\", "/"), na=False, regex=True, case=False)
        ]
        exception_path = pd.concat([exception_path, exception_pathx])
        exception_path["comment"] = "excluded path"

    scan["unc"] = scan["unc"].str.replace("\\\\", "/", regex=True)
    print("\nWARNING: files in path(s) excluded from searchx:")
    for i in paths_to_exclude:
        print("REPORT: Files to scan before applying exceptions", len(scan))
        scan = scan[
            ~scan["unc"].str.contains(
                str(i).replace("\\", "/"), na=False, regex=True, case=False
            )
        ]
        print(f"\tException: {i}")

    print("REPORT: Files to scan after applying exceptions", len(scan))

    scan["unc"] = scan["unc"].str.replace("/", "\\\\", regex=True)

    # exclude files
    # remove filenames starting with ~
    scan = scan[scan["filename"].str[0:1] != "~"]
    exception_file = scan[scan["filename"].str[0:1] == "~"]
    exception_file["comment"] = "filename starting with ~ (tilda)"

    return scan, exception_file


# %% apply_filters


def apply_filters(scan, scan_ext, scan_size):
    df_scan_ext = pd.DataFrame()
    for k, v in scan_ext.items():
        df_scan_ext_ = pd.DataFrame(v).rename(columns={0: "ext"})
        df_scan_ext_["map"] = k
        df_scan_ext = pd.concat([df_scan_ext, df_scan_ext_])

    df_scan_size = pd.DataFrame()
    for k, v in scan_size.items():
        df_scan_size_ = pd.DataFrame(v).rename(columns={0: "size_mb_"})
        df_scan_size_["map"] = k
        df_scan_size = pd.concat([df_scan_size, df_scan_size_])

    scan = pd.merge(
        scan, pd.merge(df_scan_ext, df_scan_size)[["ext", "size_mb_"]], how="left"
    )

    scan["exclude"] = np.where(scan["size_mb"] <= scan["size_mb_"], 1, 0)
    scan = scan[scan["exclude"] == 1]

    # notify if no new entries to scan
    if len(scan) == 0:
        print("REPORT: no new entries available to scan.")

    # release space
    del df_scan_ext, df_scan_size
    del scan["size_mb_"]
    del scan["exclude"]

    return scan


# %% scan drives


def scan_drives(scan, scan_ext):
    ifp = scan_folder_searchx(time_machine)
    scan = scan[~scan["uf_id"].isin(ifp["uf_id"])]

    # .txt
    load_ifp(scan, scan_ext, file_group="txt", _function=read_text)

    # .docx
    load_ifp(scan, scan_ext, file_group="docx", _function=read_docx)

    # .pptx
    load_ifp(scan, scan_ext, file_group="pptx", _function=read_pptx)

    # .msg
    load_ifp(scan, scan_ext, file_group="msg", _function=read_msg)

    # .eml
    load_ifp(scan, scan_ext, file_group="eml", _function=read_eml)

    # .epub
    load_ifp(scan, scan_ext, file_group="epub", _function=read_epub)

    # .pdf
    load_ifp(scan, scan_ext, file_group="pdf", _function=read_pdf)

    # excel
    load_ifp_xlsx(
        scan,
        scan_ext,
        file_group="excel",
        _function1=read_xlsx_text,
        _function2=read_xlsx_formula,
    )


# %% IPF (Intermediate file pool)


def ifp(scan0, searchx_final_old):
    ## load all files from ipf
    ifp_list = scan_folder_searchx(time_machine, ext=r".pickle")

    # flag files
    ifp_list["archive"] = np.where((ifp_list["uf_id"].isin(scan0["uf_id"])), 0, 1)
    keep = ifp_list[ifp_list["archive"] == 0].reset_index(drop=True)
    keep_not = ifp_list[ifp_list["archive"] == 1].reset_index(drop=True)

    ## drop old files
    print("Running time machine updates...")
    for i in range(0, len(keep_not)):
        os.remove(keep_not.loc[i, "unc"])
        print("REMOVED:", keep_not.loc[i, "unc"])

    ## load completely new files
    keep = keep[
        ~keep["uf_id"].isin(searchx_final_old["uf_id"].unique().tolist())
    ].reset_index(drop=True)
    searchx = pd.DataFrame()
    total_files = len(keep)
    print("...refreshing index with " + "{:,}".format(total_files) + " files...")

    for i in range(0, total_files):
        unc = keep.loc[i, "unc"]
        searchxx = pd.read_pickle(unc)
        # searchxx = pd.read_parquet(unc)
        searchxx = split_text(searchxx)  # limit row size
        searchxx["text"] = searchxx["text"].apply(lambda row: remove_emojis(row))
        searchxx = clean_import_load(searchxx)
        sl = len(searchxx)

        if sl != 0:
            print(
                f"READING: {i:09} of {total_files:09}; imported file length {str(sl)}"
            )
        else:
            print(
                f"READING: {i:09} of {total_files:09}; imported file length {str(sl)}; {unc}"
            )

        if len(searchxx) >= 1:
            searchxx["uf_id"] = keep.loc[i, "uf_id"]
            searchx = pd.concat([searchx, searchxx], ignore_index=True)
        else:
            pass

    ## combine old and new scan
    if len(searchx) != 0:
        searchx = pd.concat([searchx, searchx_final_old]).reset_index(drop=True)
    else:
        searchx = searchx_final_old.reset_index(drop=True)

    return searchx


def process_single_file(file_info):
    """Process a single file - designed for parallel execution"""
    try:
        unc, uf_id = file_info
        searchxx = pd.read_pickle(unc)
        searchxx = split_text(searchxx)  # limit row size
        searchxx["text"] = searchxx["text"].apply(lambda row: remove_emojis(row))
        searchxx = clean_import_load(searchxx)

        if len(searchxx) >= 1:
            searchxx["uf_id"] = uf_id
            return searchxx
        else:
            return None
    except Exception as e:
        print(f"Error processing {unc}: {e}")
        return None


def ifp_optimized(scan0, searchx_final_old, max_workers=None, batch_size=100):
    """
    Optimized version with parallel processing and batching

    Args:
        scan0: DataFrame with uf_id column
        searchx_final_old: DataFrame with existing data
        max_workers: Number of parallel workers (default: CPU count)
        batch_size: Number of files to process in each batch
    """

    # Set default workers to CPU count
    if max_workers is None:
        max_workers = min(mp.cpu_count(), 8)  # Cap at 8 to avoid too many processes

    print(f"Using {max_workers} parallel workers")

    ## Load all files from ifp
    ifp_list = scan_folder_searchx(time_machine, ext=r".pickle")

    # Convert to sets for faster lookup
    scan0_ids = set(scan0["uf_id"].tolist())
    existing_ids = set(searchx_final_old["uf_id"].unique().tolist())

    # Flag files using vectorized operations
    ifp_list["archive"] = np.where(ifp_list["uf_id"].isin(scan0_ids), 0, 1)

    # Split into keep and remove
    keep = ifp_list[ifp_list["archive"] == 0].copy()
    keep_not = ifp_list[ifp_list["archive"] == 1].copy()

    ## Remove old files in parallel (I/O bound - use ThreadPoolExecutor)
    if len(keep_not) > 0:
        print("Running time machine updates...")
        files_to_remove = keep_not["unc"].tolist()

        def remove_file(filepath):
            try:
                os.remove(filepath)
                return f"REMOVED: {filepath}"
            except Exception as e:
                return f"ERROR removing {filepath}: {e}"

        with ThreadPoolExecutor(max_workers=min(max_workers, 4)) as executor:
            removal_results = list(executor.map(remove_file, files_to_remove))
            for result in removal_results:
                print(result)

    ## Filter to completely new files using set operations
    keep = keep[~keep["uf_id"].isin(existing_ids)].reset_index(drop=True)

    total_files = len(keep)
    print(f"...refreshing index with {total_files:,} files...")

    if total_files == 0:
        return searchx_final_old.reset_index(drop=True)

    # Prepare file info for parallel processing
    file_info_list = [(row["unc"], row["uf_id"]) for _, row in keep.iterrows()]

    # Process files in batches to manage memory
    all_results = []

    for batch_start in range(0, total_files, batch_size):
        batch_end = min(batch_start + batch_size, total_files)
        batch_files = file_info_list[batch_start:batch_end]

        print(
            f"Processing batch {batch_start // batch_size + 1} of {(total_files - 1) // batch_size + 1} "
            f"(files {batch_start + 1}-{batch_end})"
        )

        # Process batch in parallel
        with ProcessPoolExecutor(max_workers=max_workers) as executor:
            batch_results = list(executor.map(process_single_file, batch_files))

        # Filter out None results and collect valid DataFrames
        valid_results = [df for df in batch_results if df is not None and len(df) >= 1]

        if valid_results:
            # Concatenate batch results
            batch_df = pd.concat(valid_results, ignore_index=True)
            all_results.append(batch_df)
            print(f"Batch completed: {len(valid_results)} valid files processed")

    # Combine all results
    if all_results:
        searchx = pd.concat(all_results, ignore_index=True)
        # Combine with existing data
        searchx = pd.concat([searchx, searchx_final_old], ignore_index=True)
    else:
        searchx = searchx_final_old.reset_index(drop=True)

    return searchx


# Alternative: Memory-efficient version that processes files one by one but with optimizations
def ifp_memory_efficient(scan0, searchx_final_old, progress_interval=50):
    """
    Memory-efficient version that processes files sequentially but with optimizations
    Better for systems with limited RAM
    """

    ## Load all files from ifp
    ifp_list = scan_folder_searchx(time_machine, ext=r".pickle")

    # Convert to sets for faster lookup
    scan0_ids = set(scan0["uf_id"].tolist())
    existing_ids = set(searchx_final_old["uf_id"].unique().tolist())

    # Vectorized operations
    ifp_list["archive"] = np.where(ifp_list["uf_id"].isin(scan0_ids), 0, 1)
    keep = ifp_list[ifp_list["archive"] == 0].copy()
    keep_not = ifp_list[ifp_list["archive"] == 1].copy()

    ## Remove old files
    if len(keep_not) > 0:
        print("Running time machine updates...")
        for _, row in keep_not.iterrows():
            try:
                os.remove(row["unc"])
                print(f"REMOVED: {row['unc']}")
            except Exception as e:
                print(f"ERROR removing {row['unc']}: {e}")

    ## Filter to new files
    keep = keep[~keep["uf_id"].isin(existing_ids)].reset_index(drop=True)

    total_files = len(keep)
    print(f"...refreshing index with {total_files:,} files...")

    if total_files == 0:
        return searchx_final_old.reset_index(drop=True)

    # Process files with periodic concatenation to manage memory
    searchx_parts = []
    current_batch = []

    for i in range(total_files):
        try:
            unc = keep.loc[i, "unc"]
            uf_id = keep.loc[i, "uf_id"]

            searchxx = pd.read_pickle(unc)
            searchxx = split_text(searchxx)
            searchxx["text"] = searchxx["text"].apply(lambda row: remove_emojis(row))
            searchxx = clean_import_load(searchxx)

            sl = len(searchxx)

            if i % progress_interval == 0 or i == total_files - 1:
                print(
                    f"READING: {i + 1:09} of {total_files:09}; imported file length {sl}"
                )

            if sl >= 1:
                searchxx["uf_id"] = uf_id
                current_batch.append(searchxx)

                # Concatenate every 100 files to manage memory
                if len(current_batch) >= 100:
                    batch_df = pd.concat(current_batch, ignore_index=True)
                    searchx_parts.append(batch_df)
                    current_batch = []

        except Exception as e:
            print(f"Error processing file {i}: {e}")
            continue

    # Handle remaining files in current_batch
    if current_batch:
        batch_df = pd.concat(current_batch, ignore_index=True)
        searchx_parts.append(batch_df)

    # Final concatenation
    if searchx_parts:
        searchx = pd.concat(searchx_parts, ignore_index=True)
        searchx = pd.concat([searchx, searchx_final_old], ignore_index=True)
    else:
        searchx = searchx_final_old.reset_index(drop=True)

    return searchx


# %% Export indexes


def export_index_files(_fs_index_dir, _time_machine_path, scan0, searchx):
    print("REPORT: Exporting final outputs...")
    os.chdir(_fs_index_dir)

    # convert to polars dataframe
    try:
        scan0_polars = pl.DataFrame(scan0.replace("", None))
        searchx_polars = pl.DataFrame(searchx.replace("", None))
    except Exception:
        scan0_polars = pl.DataFrame(scan0)
        searchx_polars = pl.DataFrame(searchx)

    # parquet: issues reading writing both with polars and pandas
    try:
        scan0_polars.write_parquet("scan0.parquet")
        searchx_polars.write_parquet("searchx_final.parquet")
        print(f"REPORT: searchx, row_count = {len(searchx):,}")
    except Exception:
        print("ERROR: Issues exporting as paraquet.")

    # zip time machine
    zip_folder(_time_machine_path, str(_time_machine_path) + ".zip")

    # delete time machine folder
    delete_folder(_time_machine_path)

    # delete scan.csv
    # os.remove(_fs_index_dir / "scan.csv")

    # %% end monitoring progress

    timer_end()


# %% scan_directories_python


def format_datetime(timestamp):
    """Convert timestamp to formatted datetime string"""
    if timestamp:
        return datetime.fromtimestamp(timestamp).strftime("%m/%d/%Y %I:%M:%S %p")
    return ""


def format_datetime_utc(timestamp):
    """Convert timestamp to UTC formatted datetime string"""
    if timestamp:
        # Create a UTC timezone object
        utc_timezone = timezone.utc
        return datetime.fromtimestamp(timestamp, tz=utc_timezone).strftime(
            "%m/%d/%Y %I:%M:%S %p"
        )
    return ""


def scan_directories_python(fs_index_dir, scan_dirs):
    """
    Scan directories and collect file information directly in Python

    Args:
        fs_index_dir (str): Directory where the output CSV will be saved
        scan_dirs (list): List of directories to scan
    """

    output_file = os.path.join(fs_index_dir, "scan.csv")

    # Ensure output directory exists
    os.makedirs(fs_index_dir, exist_ok=True)

    print("Starting Python-based directory scan...")
    # print(f"Output file: {output_file}")

    # Open CSV file for writing
    with open(output_file, "w", newline="", encoding="utf-8") as csvfile:
        # Define CSV headers to match PowerShell output
        fieldnames = [
            "Owner",
            "FullName",
            "PSChildName",
            "CreationTime",
            "CreationTimeUtc",
            "LastAccessTime",
            "LastAccessTimeUtc",
            "LastWriteTime",
            "LastWriteTimeUtc",
            "Length",
            "Extension",
        ]

        writer = csv.DictWriter(csvfile, fieldnames=fieldnames, delimiter="|")
        writer.writeheader()

        total_files = 0

        for scan_dir in scan_dirs:
            print(f"Scanning directory: {scan_dir}")

            if not os.path.exists(scan_dir):
                print(f"WARNING: Directory does not exist: {scan_dir}")
                continue

            dir_file_count = 0

            try:
                # Walk through directory recursively
                for root, dirs, files in os.walk(scan_dir):
                    for file in files:
                        try:
                            file_path = os.path.join(root, file)

                            # Get file statistics
                            file_stat = os.stat(file_path)

                            # Get file extension
                            _, extension = os.path.splitext(file)

                            # Get file owner
                            # owner = get_file_owner(file_path)

                            # Prepare file information
                            file_info = {
                                "Owner": "",
                                "FullName": file_path,
                                "PSChildName": file,
                                "CreationTime": format_datetime(file_stat.st_ctime),
                                "CreationTimeUtc": format_datetime_utc(
                                    file_stat.st_ctime
                                ),
                                "LastAccessTime": format_datetime(file_stat.st_atime),
                                "LastAccessTimeUtc": format_datetime_utc(
                                    file_stat.st_atime
                                ),
                                "LastWriteTime": format_datetime(file_stat.st_mtime),
                                "LastWriteTimeUtc": format_datetime_utc(
                                    file_stat.st_mtime
                                ),
                                "Length": file_stat.st_size,
                                "Extension": extension,
                            }

                            # Write to CSV
                            writer.writerow(file_info)
                            dir_file_count += 1
                            total_files += 1

                            # Progress indicator
                            if dir_file_count % 1000 == 0:
                                print(
                                    f"  Processed {dir_file_count} files in {scan_dir}..."
                                )

                        except (OSError, PermissionError) as e:
                            print(f"  WARNING: Could not access file {file_path}: {e}")
                            continue

            except (OSError, PermissionError) as e:
                print(f"ERROR: Could not access directory {scan_dir}: {e}")
                continue

            print(
                rf"REPORT: Completed scanning {scan_dir}: {dir_file_count} files processed"
            )

    print(f"Total files processed: {total_files}")
    print(f"Output saved to: {output_file}")


# Alternative function that can be used as a drop-in replacement
def scan_python_replacement(fs_index_dir, scan_dirs):
    """
    Drop-in replacement for the original scan_powershell function
    """
    print("\nNOTE: Using Python-based scanning instead of PowerShell...")
    scan_directories_python(fs_index_dir, scan_dirs)
