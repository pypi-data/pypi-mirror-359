import os
import random
import pandas as pd
from faker import Faker
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. PowerPoint files will not be generated.")

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("Warning: reportlab not installed. PDF files will not be generated.")


class FakeDataGenerator:
    def __init__(self, keywords=None):
        self.fake = Faker()
        if keywords is not None:
            self.keywords = keywords
        else:
            self.keywords = [
                "key attributes",
                "random forests",
                "class",
                "disease_type",
                "color_hex",
                "raw_score_ocs_pre",
                "predicted values",
                "element",
                "hcp_id",
                "sell price",
                "stage section",
                "s&p 500",
                "tbl",
                "field",
                "priority",
                "weights",
                "fill5",
                "company",
                "derived_in",
                "tutorial",
                "chapter three",
                "fill_date",
                "command",
                "sort_order",
                "mode_of_dispensation",
                "groupb",
                "for linear data",
                "revenue",
                "end_date",
                "chart_3",
                "varnum",
                "ocs_top50",
                "ocs_months_pre",
                "remaining",
                "response",
                "category",
                "k samples",
                "drug_name",
                "chapter eleven",
                "total_months_post",
                "% variation",
            ]

    def get_short_keywords(self, max_length=10):
        """Get keywords shorter than specified length"""
        return [kw for kw in self.keywords if len(kw) < max_length]

    def get_very_short_keywords(self, max_length=6):
        """Get very short keywords for table names"""
        return [kw for kw in self.keywords if len(kw) < max_length]

    def generate_excel_files(self, num_files=5, max_rows=100, output_dir="fake_data"):
        """Generate Excel files (.xlsx and .xls) with keywords as column names"""
        short_keywords = self.get_short_keywords()

        for i in range(num_files):
            # Use .xlsx format (more modern and widely supported)
            extension = ".xlsx"
            filename = f"data_{i + 1}{extension}"
            filepath = os.path.join(output_dir, filename)

            # Random number of columns (2-15)
            num_cols = random.randint(2, 15)
            columns = random.sample(short_keywords, min(num_cols, len(short_keywords)))

            # Random number of rows
            num_rows = random.randint(10, max_rows)

            # Generate data
            data = {}
            for col in columns:
                col_data = []
                for _ in range(num_rows):
                    # Mix of different data types and include keywords
                    if random.random() < 0.3:  # 30% chance of using a keyword
                        col_data.append(random.choice(self.keywords))
                    elif random.random() < 0.5:  # 50% chance of text
                        col_data.append(self.fake.sentence())
                    else:  # Numbers
                        col_data.append(random.randint(1, 10000))
                data[col] = col_data

            df = pd.DataFrame(data)

            # Generate .xlsx files with openpyxl engine
            try:
                df.to_excel(filepath, index=False, engine="openpyxl")
                print(f"Generated: {filename}")
            except Exception:
                pass

    def generate_pptx_files(self, num_files=3, max_slides=10, output_dir="fake_data"):
        """Generate PowerPoint files with keywords embedded in slides"""
        if not PPTX_AVAILABLE:
            print("Skipping PowerPoint generation - python-pptx not installed")
            return

        for i in range(num_files):
            filename = f"presentation_{i + 1}.pptx"
            filepath = os.path.join(output_dir, filename)

            # Create presentation
            prs = Presentation()

            # Random number of slides
            num_slides = random.randint(3, max_slides)

            for slide_num in range(num_slides):
                # Add slide with title and content layout
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)

                # Title with keywords
                title = slide.shapes.title
                title_text = random.choice(self.keywords)
                if len(title_text) > 50:  # Shorten long titles
                    title_text = title_text[:47] + "..."
                title.text = f"Slide {slide_num + 1}: {title_text}"

                # Content with keywords
                content = slide.placeholders[1]
                text_frame = content.text_frame
                text_frame.clear()

                # Add multiple bullet points with keywords
                num_bullets = random.randint(2, 6)
                for bullet_num in range(num_bullets):
                    p = (
                        text_frame.paragraphs[0]
                        if bullet_num == 0
                        else text_frame.add_paragraph()
                    )

                    # Mix keywords with fake text
                    if random.random() < 0.6:  # 60% chance of using keywords
                        bullet_text = (
                            f"{random.choice(self.keywords)} - {self.fake.sentence()}"
                        )
                    else:
                        bullet_text = self.fake.sentence()

                    p.text = bullet_text
                    p.level = 0

            # Save presentation
            try:
                prs.save(filepath)
                print(f"Generated: {filename}")
            except Exception as e:
                print(f"Error generating {filename}: {e}")

    def generate_pdf_files(self, num_files=3, max_pages=5, output_dir="fake_data"):
        """Generate PDF files with keywords embedded in text"""
        if not PDF_AVAILABLE:
            print("Skipping PDF generation - reportlab not installed")
            return

        styles = getSampleStyleSheet() if PDF_AVAILABLE else None

        for i in range(num_files):
            filename = f"document_{i + 1}.pdf"
            filepath = os.path.join(output_dir, filename)

            try:
                # Create PDF document
                doc = SimpleDocTemplate(filepath, pagesize=letter)
                story = []

                # Title
                title_keyword = random.choice(self.keywords)
                if len(title_keyword) > 60:
                    title_keyword = title_keyword[:57] + "..."
                title = Paragraph(f"<b>{title_keyword}</b>", styles["Title"])
                story.append(title)
                story.append(Spacer(1, 12))

                # Generate content for multiple pages
                num_pages = random.randint(2, max_pages)

                for page_num in range(num_pages):
                    # Page heading
                    page_title = (
                        f"Section {page_num + 1}: {random.choice(self.keywords)}"
                    )
                    if len(page_title) > 80:
                        page_title = page_title[:77] + "..."

                    heading = Paragraph(f"<b>{page_title}</b>", styles["Heading1"])
                    story.append(heading)
                    story.append(Spacer(1, 12))

                    # Multiple paragraphs per page
                    num_paragraphs = random.randint(3, 6)

                    for para_num in range(num_paragraphs):
                        # Create paragraph with keywords
                        if random.random() < 0.5:  # 50% chance of starting with keyword
                            para_text = f"{random.choice(self.keywords)} {self.fake.paragraph()}"
                        else:
                            para_text = f"{self.fake.paragraph()} {random.choice(self.keywords)}"

                        # Add more keywords randomly within the paragraph
                        if random.random() < 0.3:
                            para_text += f" {random.choice(self.keywords)}"

                        paragraph = Paragraph(para_text, styles["Normal"])
                        story.append(paragraph)
                        story.append(Spacer(1, 12))

                    # Add space between sections
                    story.append(Spacer(1, 20))

                # Build PDF
                doc.build(story)
                print(f"Generated: {filename}")

            except Exception as e:
                print(f"Error generating {filename}: {e}")
            except Exception as e:
                print(f"Error generating {filename}: {e}")

    def generate_txt_files(self, num_files=5, max_lines=50, output_dir="fake_data"):
        """Generate text files with keywords embedded in random text"""
        for i in range(num_files):
            filename = f"text_{i + 1}.txt"
            filepath = os.path.join(output_dir, filename)

            num_lines = random.randint(10, max_lines)

            with open(filepath, "w", encoding="utf-8") as f:
                for _ in range(num_lines):
                    # Mix regular sentences with keyword sentences
                    if random.random() < 0.4:  # 40% chance of using keywords
                        keyword = random.choice(self.keywords)
                        sentence = (
                            f"{self.fake.sentence()} {keyword} {self.fake.sentence()}"
                        )
                    else:
                        sentence = self.fake.paragraph(
                            nb_sentences=random.randint(1, 3)
                        )
                    f.write(sentence + "\n")

            print(f"Generated: {filename}")

    def generate_sql_files(self, num_files=5, max_queries=20, output_dir="fake_data"):
        """Generate SQL files with keywords as column and table names"""
        short_keywords = self.get_short_keywords()
        very_short_keywords = self.get_very_short_keywords()

        for i in range(num_files):
            filename = f"queries_{i + 1}.sql"
            filepath = os.path.join(output_dir, filename)

            num_queries = random.randint(5, max_queries)

            with open(filepath, "w", encoding="utf-8") as f:
                f.write(f"-- SQL File {i + 1} Generated on {datetime.now()}\n")
                f.write("-- Contains keywords for text mining purposes\n\n")

                for j in range(num_queries):
                    # Choose random table name from very short keywords
                    table_name = (
                        random.choice(very_short_keywords)
                        if very_short_keywords
                        else "tbl"
                    )
                    table_name = table_name.replace(" ", "_").replace("%", "pct")

                    # Choose random columns from short keywords
                    num_cols = random.randint(1, 4)
                    columns = random.sample(
                        short_keywords, min(num_cols, len(short_keywords))
                    )
                    column_names = [
                        col.replace(" ", "_").replace("%", "pct") for col in columns
                    ]

                    # Generate different types of SQL queries
                    query_type = random.choice(["SELECT", "INSERT", "UPDATE", "CREATE"])

                    if query_type == "SELECT":
                        f.write(f"SELECT {', '.join(column_names)}\n")
                        f.write(f"FROM {table_name}\n")
                        if random.random() < 0.5:
                            f.write(
                                f"WHERE {column_names[0]} = '{random.choice(self.keywords)}'\n"
                            )
                        f.write(";\n\n")

                    elif query_type == "CREATE":
                        f.write(f"CREATE TABLE {table_name} (\n")
                        for k, col in enumerate(column_names):
                            data_type = random.choice(
                                ["VARCHAR(255)", "INT", "DECIMAL(10,2)", "DATE"]
                            )
                            comma = "," if k < len(column_names) - 1 else ""
                            f.write(f"    {col} {data_type}{comma}\n")
                        f.write(");\n\n")

                    elif query_type == "INSERT":
                        f.write(
                            f"INSERT INTO {table_name} ({', '.join(column_names)})\n"
                        )
                        f.write("VALUES (")
                        values = []
                        for col in column_names:
                            if random.random() < 0.3:
                                values.append(f"'{random.choice(self.keywords)}'")
                            else:
                                values.append(f"'{self.fake.word()}'")
                        f.write(", ".join(values))
                        f.write(");\n\n")

            print(f"Generated: {filename}")

    def generate_python_files(self, num_files=5, max_lines=30, output_dir="fake_data"):
        """Generate Python files with keywords as variables and in comments"""
        short_keywords = self.get_short_keywords()
        very_short_keywords = self.get_very_short_keywords()

        for i in range(num_files):
            filename = f"script_{i + 1}.py"
            filepath = os.path.join(output_dir, filename)

            with open(filepath, "w", encoding="utf-8") as f:
                f.write(f"#!/usr/bin/env python3\n")
                f.write(f"# Python script {i + 1} - Generated for text mining\n")
                f.write(
                    f"# Contains keywords: {', '.join(random.sample(self.keywords, 3))}\n\n"
                )

                f.write(
                    "import random\nimport pandas as pd\nfrom datetime import datetime\n\n"
                )

                # Generate some functions with keyword-based names and content
                num_functions = random.randint(2, 4)

                for j in range(num_functions):
                    # Function name from keywords
                    func_keyword = (
                        random.choice(very_short_keywords)
                        if very_short_keywords
                        else "func"
                    )
                    func_name = (
                        func_keyword.replace(" ", "_").replace("%", "pct").lower()
                    )

                    f.write(f"def {func_name}():\n")
                    f.write(f'    """{random.choice(self.keywords)}"""\n')

                    # Generate function body with keywords
                    num_lines = random.randint(3, 8)
                    for k in range(num_lines):
                        if random.random() < 0.4:
                            var_name = (
                                random.choice(short_keywords)
                                .replace(" ", "_")
                                .replace("%", "pct")
                                .lower()
                            )
                            f.write(
                                f"    {var_name} = '{random.choice(self.keywords)}'\n"
                            )
                        else:
                            f.write(f"    # {random.choice(self.keywords)}\n")
                            f.write(f"    result = '{self.fake.sentence()}'\n")

                    f.write("    return result\n\n")

                # Main execution block
                f.write('if __name__ == "__main__":\n')
                f.write(f'    print("{random.choice(self.keywords)}")\n')
                for keyword in random.sample(
                    short_keywords, min(3, len(short_keywords))
                ):
                    var_name = keyword.replace(" ", "_").replace("%", "pct").lower()
                    f.write(f'    {var_name} = "{random.choice(self.keywords)}"\n')

            print(f"Generated: {filename}")

    def generate_all_files(
        self,
        xlsx_count=3,
        txt_count=3,
        sql_count=3,
        py_count=3,
        pptx_count=2,
        pdf_count=2,
        max_rows=100,
        max_lines=50,
        max_slides=10,
        max_pages=5,
        output_dir="fake_data",
    ):
        """Generate all types of files"""
        # Create output directory
        os.makedirs(output_dir, exist_ok=True)

        print(f"Generating fake data files in '{output_dir}' directory...")
        print(f"Using {len(self.keywords)} keywords for text mining corpus\n")

        # Generate each file type
        if xlsx_count > 0:
            print("Generating Excel files...")
            self.generate_excel_files(xlsx_count, max_rows, output_dir)
            print()

        if txt_count > 0:
            print("Generating text files...")
            self.generate_txt_files(txt_count, max_lines, output_dir)
            print()

        if sql_count > 0:
            print("Generating SQL files...")
            self.generate_sql_files(sql_count, max_lines, output_dir)
            print()

        if py_count > 0:
            print("Generating Python files...")
            self.generate_python_files(py_count, max_lines, output_dir)
            print()

        if pptx_count > 0:
            print("Generating PowerPoint files...")
            self.generate_pptx_files(pptx_count, max_slides, output_dir)
            print()

        if pdf_count > 0:
            print("Generating PDF files...")
            self.generate_pdf_files(pdf_count, max_pages, output_dir)
            print()

        print(f"✅ All files generated successfully in '{output_dir}' directory!")

    def generate_fake_info_schema(self, num_rows=100):
        """
        Generate a fake information schema DataFrame with keywords as values

        Args:
            num_rows (int): Number of rows to generate in the DataFrame

        Returns:
            pandas.DataFrame: DataFrame with columns ['table_catalog', 'table_schema', 'table_name', 'column_name']
        """
        # Filter keywords to those with length <= 10 characters
        short_keywords = [kw for kw in self.keywords if len(kw) <= 10]

        # If no keywords are short enough, use some default values
        if not short_keywords:
            short_keywords = ["data", "info", "test", "main", "temp", "prod", "dev"]

        # Generate sample catalogs, schemas - typically shorter names
        catalogs = ["main", "test", "prod", "dev", "staging"]
        schemas = ["public", "dbo", "schema1", "data", "analytics", "reports"]

        # Ensure we have enough variety, add some from short keywords
        catalogs.extend(
            [
                kw.replace(" ", "_").replace("%", "pct").lower()
                for kw in short_keywords[:5]
            ]
        )
        schemas.extend(
            [
                kw.replace(" ", "_").replace("%", "pct").lower()
                for kw in short_keywords[:8]
            ]
        )

        # Generate table names from keywords (clean them for SQL compatibility)
        table_names = [
            kw.replace(" ", "_").replace("%", "pct").lower() for kw in short_keywords
        ]

        # Generate column names from keywords (clean them for SQL compatibility)
        column_names = [
            kw.replace(" ", "_").replace("%", "pct").lower() for kw in short_keywords
        ]

        # Generate the DataFrame
        data = []
        for _ in range(num_rows):
            row = {
                "table_catalog": random.choice(catalogs),
                "table_schema": random.choice(schemas),
                "table_name": random.choice(table_names),
                "column_name": random.choice(column_names),
            }
            data.append(row)

        df = pd.DataFrame(data)

        print(f"Generated fake information schema with {num_rows} rows")
        print(f"Unique catalogs: {df['table_catalog'].nunique()}")
        print(f"Unique schemas: {df['table_schema'].nunique()}")
        print(f"Unique tables: {df['table_name'].nunique()}")
        print(f"Unique columns: {df['column_name'].nunique()}")

        return df


# Example usage
if __name__ == "__main__":
    generator = FakeDataGenerator()

    # Generate a fake information schema with 50 rows
    info_schema_df = generator.generate_fake_info_schema(num_rows=50)

    # Display first few rows
    print("\nFirst 10 rows of fake information schema:")
    print(info_schema_df.head(10))

    # Save to CSV if needed
    info_schema_df.to_csv("fake_info_schema.csv", index=False)
    print("\nSaved to 'fake_info_schema.csv'")

    # Generate a larger one
    big_schema_df = generator.generate_fake_info_schema(num_rows=500)

    # Show some statistics
    print("\nSample of unique values:")
    print(f"Catalogs: {list(big_schema_df['table_catalog'].unique())[:5]}")
    print(f"Schemas: {list(big_schema_df['table_schema'].unique())[:5]}")
    print(f"Tables: {list(big_schema_df['table_name'].unique())[:5]}")
    print(f"Columns: {list(big_schema_df['column_name'].unique())[:5]}")


# Example usage
if __name__ == "__main__":
    generator = FakeDataGenerator()

    # Generate files with custom parameters
    generator.generate_all_files(
        xlsx_count=4,  # Number of Excel files (.xlsx)
        txt_count=3,  # Number of text files
        sql_count=2,  # Number of SQL files
        py_count=2,  # Number of Python files
        pptx_count=3,  # Number of PowerPoint files
        pdf_count=2,  # Number of PDF files
        max_rows=50,  # Maximum rows for Excel files
        max_lines=30,  # Maximum lines for text/SQL/Python files
        max_slides=8,  # Maximum slides for PowerPoint files
        max_pages=4,  # Maximum pages for PDF files
        output_dir="fake_data_corpus",  # Output directory
    )

    # You can also generate specific file types individually:
    # generator.generate_excel_files(num_files=5, max_rows=100)
    # generator.generate_txt_files(num_files=3, max_lines=25)
    # generator.generate_sql_files(num_files=2, max_queries=15)
    # generator.generate_python_files(num_files=2, max_lines=20)
    # generator.generate_pptx_files(num_files=3, max_slides=10)
    # generator.generate_pdf_files(num_files=2, max_pages=5)
