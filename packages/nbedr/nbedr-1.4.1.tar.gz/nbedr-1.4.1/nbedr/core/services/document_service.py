"""
Document processing service for chunking, text extraction, and embedding generation.
"""

import fcntl
import hashlib
import json
import logging
import re
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from math import ceil
from pathlib import Path
from types import ModuleType
from typing import IO, Any, Callable, Dict, Iterator, List, Optional, Set, Type, TypeVar, Union, cast

from ..config import EmbeddingConfig
from ..models import ChunkingStrategy, DocType, DocumentChunk, EmbeddingBatch, EmbeddingStats, ProcessingResult
from ..utils.document_coordinator import DocumentCoordinator
from ..utils.instance_coordinator import InstanceCoordinator, create_instance_coordinator

# Type variables for generic types
T = TypeVar("T")

try:
    import pypdf

    PypdfModule = Type[ModuleType]
except ImportError:
    pypdf = None
    PypdfModule = Type[Any]  # type: ignore[misc]

try:
    from pptx import Presentation

    PresentationCallable = Type[Callable[[Union[str, IO[bytes], None]], Any]]
except ImportError:
    # Create mock class for when python-pptx is not available
    class _MockPresentation:
        pass

    Presentation = _MockPresentation  # type: ignore
    PresentationCallable = Type[Any]  # type: ignore[misc]

try:
    from langchain_experimental.text_splitter import SemanticChunker
    from langchain_openai.embeddings import AzureOpenAIEmbeddings, OpenAIEmbeddings

    SemanticChunkerType = Type[SemanticChunker]
    OpenAIEmbeddingsType = Type[OpenAIEmbeddings]
    AzureOpenAIEmbeddingsType = Type[AzureOpenAIEmbeddings]
except ImportError:
    # Create mock classes for when langchain is not available
    class _MockSemanticChunker:
        pass

    class _MockOpenAIEmbeddings:
        pass

    class _MockAzureOpenAIEmbeddings:
        pass

    SemanticChunker = _MockSemanticChunker  # type: ignore
    OpenAIEmbeddings = _MockOpenAIEmbeddings  # type: ignore
    AzureOpenAIEmbeddings = _MockAzureOpenAIEmbeddings  # type: ignore
    SemanticChunkerType = Type[Any]  # type: ignore[misc]
    OpenAIEmbeddingsType = Type[Any]  # type: ignore[misc]
    AzureOpenAIEmbeddingsType = Type[Any]  # type: ignore[misc]

try:
    from tqdm.auto import tqdm as tqdm_auto

    def tqdm(iterable: Iterator[T], *args: Any, **kwargs: Any) -> Iterator[T]:
        """Wrapper around tqdm to handle both CLI and notebook contexts."""
        return cast(Iterator[T], tqdm_auto(iterable, *args, **kwargs))

except ImportError:

    def tqdm(iterable: Iterator[T], *args: Any, **kwargs: Any) -> Iterator[T]:
        """Fallback implementation when tqdm is not available."""
        return iterable


logger = logging.getLogger(__name__)


class DocumentService:
    """Service for document processing, chunking, and embedding generation."""

    def __init__(self, config: EmbeddingConfig, enable_coordination: bool = True):
        self.config = config
        self.stats = EmbeddingStats()

        # Initialize instance coordination if enabled
        self.coordinator = None
        self.document_coordinator = None
        if enable_coordination:
            try:
                self.coordinator = create_instance_coordinator(config)
                if self.coordinator.register_instance():
                    logger.info("Instance coordination enabled")

                    # Initialize document coordination
                    self.document_coordinator = DocumentCoordinator(
                        coordination_dir=self.coordinator.coordination_dir, instance_id=self.coordinator.instance_id
                    )

                    # Update config with instance-specific paths if suggested
                    suggested_paths = self.coordinator.suggest_instance_specific_paths()
                    if "output_path" in suggested_paths:
                        self.config.output = suggested_paths["output_path"]
                        logger.info(f"Using instance-specific output path: {self.config.output}")
                else:
                    logger.error("Failed to register instance - conflicts detected")
                    self.coordinator = None
            except Exception as e:
                logger.warning(f"Failed to initialize instance coordination: {e}")
                self.coordinator = None

    def process_documents(self, data_path: Path) -> List[DocumentChunk]:
        """Process documents and return chunks."""
        logger.info(f"Processing documents from {data_path}")

        if self.config.doctype == "api":
            return self._process_api_documents(data_path)
        else:
            return self._process_regular_documents(data_path)

    def generate_embeddings(self, chunks: List[DocumentChunk]) -> List[DocumentChunk]:
        """Generate embeddings for document chunks."""
        logger.info(f"Generating embeddings for {len(chunks)} chunks")

        # TODO: Implement actual embedding generation
        # This is where you would integrate with OpenAI API or other embedding services
        embedded_chunks = []

        # Process in batches
        batch_size = self.config.batch_size_embeddings
        for i in range(0, len(chunks), batch_size):
            batch_chunks = chunks[i : i + batch_size]
            batch = EmbeddingBatch.create(batch_chunks, self.config.embedding_model)

            try:
                processed_batch = self._process_embedding_batch(batch)
                embedded_chunks.extend(processed_batch.chunks)
                batch.mark_completed()
                logger.info(f"Processed batch {i // batch_size + 1}/{ceil(len(chunks) / batch_size)}")
            except Exception as e:
                logger.error(f"Failed to process embedding batch: {e}")
                batch.mark_failed(str(e))
                # Continue with remaining batches

        return embedded_chunks

    def _process_embedding_batch(self, batch: EmbeddingBatch) -> EmbeddingBatch:
        """Process a batch of chunks for embedding generation."""
        # TODO: Implement actual embedding API calls
        # For now, return mock embeddings

        texts = [chunk.content for chunk in batch.chunks]

        # Mock embedding generation - replace with real implementation
        mock_embeddings = self._generate_mock_embeddings(texts)

        # Update chunks with embeddings
        for chunk, embedding in zip(batch.chunks, mock_embeddings):
            chunk.embedding = embedding
            chunk.embedding_model = batch.model

        return batch

    def _generate_mock_embeddings(self, texts: List[str]) -> List[List[float]]:
        """Generate mock embeddings for testing purposes."""
        # TODO: Replace with actual embedding generation
        # This would call OpenAI API, Azure OpenAI, or other embedding services

        import random

        dimension = self.config.embedding_dimensions

        embeddings = []
        for text in texts:
            # Generate deterministic mock embedding based on text hash
            # Note: This is for testing/mocking only, not cryptographic use
            random.seed(hash(text) % (2**32))  # nosec B311
            embedding = [random.uniform(-1, 1) for _ in range(dimension)]  # nosec B311
            embeddings.append(embedding)

        return embeddings

    def store_embeddings(self, chunks: List[DocumentChunk]) -> bool:
        """Store embeddings in vector database."""
        logger.info(f"Storing {len(chunks)} embeddings in {self.config.vector_db_type} database")

        # TODO: Implement vector database integration
        # This would integrate with FAISS, Pinecone, Chroma, etc.

        try:
            # Mock storage operation
            for chunk in chunks:
                if chunk.has_embedding():
                    # Generate mock vector ID
                    chunk.vector_id = f"vec_{chunk.id}"

            logger.info(f"Successfully stored {len(chunks)} embeddings")
            return True

        except Exception as e:
            logger.error(f"Failed to store embeddings: {e}")
            return False

    def _process_api_documents(self, data_path: Path) -> List[DocumentChunk]:
        """Process API documentation from JSON file."""
        with open(data_path) as f:
            api_docs_json = json.load(f)

        required_fields = ["user_name", "api_name", "api_call", "api_version", "api_arguments", "functionality"]
        if api_docs_json and isinstance(api_docs_json[0], dict):
            for field in required_fields:
                if field not in api_docs_json[0]:
                    raise ValueError(f"API documentation missing required field: {field}")

        chunks = []
        for i, api_doc in enumerate(api_docs_json):
            chunk = DocumentChunk.create(
                content=str(api_doc), source=str(data_path), metadata={"type": "api", "index": i}
            )
            chunks.append(chunk)

        self.stats.total_chunks = len(chunks)
        return chunks

    def _process_regular_documents(self, data_path: Path) -> List[DocumentChunk]:
        """Process regular documents (PDF, TXT, JSON, PPTX) with contention prevention."""
        # Get list of files to process
        file_paths = []
        if data_path.is_dir():
            file_paths = list(data_path.rglob(f"**/*.{self.config.doctype}"))
        else:
            file_paths = [data_path]

        # Filter out files that are already being processed or completed
        available_files = self._get_available_files(file_paths)

        if not available_files:
            logger.info("No available files to process (all may be in progress or completed)")
            return []

        logger.info(f"Processing {len(available_files)} available files out of {len(file_paths)} total")

        all_chunks = []
        futures = []

        # Use a simple progress counter instead of tqdm
        total_files = len(available_files)
        processed_files = 0

        with ThreadPoolExecutor(max_workers=self.config.embed_workers) as executor:
            for file_path in available_files:
                future = executor.submit(self._process_single_file_with_coordination, file_path)
                futures.append(future)

                if self.config.pace:
                    time.sleep(15)

            for future in as_completed(futures):
                try:
                    chunks = future.result()
                    if chunks:  # Only extend if chunks were successfully processed
                        all_chunks.extend(chunks)
                    processed_files += 1
                    logger.info(f"Processed {processed_files}/{total_files} files, total chunks: {len(all_chunks)}")
                except Exception as e:
                    logger.error(f"Error processing file: {e}")
                    processed_files += 1

        self.stats.total_chunks = len(all_chunks)
        return all_chunks

    def update_heartbeat(self):
        """Update instance heartbeat if coordination is enabled."""
        if self.coordinator:
            try:
                self.coordinator.update_heartbeat()
            except Exception as e:
                logger.warning(f"Failed to update heartbeat: {e}")

    def cleanup_instance(self):
        """Cleanup instance registration."""
        if self.coordinator:
            try:
                self.coordinator.unregister_instance()
                logger.info("Instance coordination cleaned up")
            except Exception as e:
                logger.warning(f"Failed to cleanup instance coordination: {e}")

    def __del__(self):
        """Cleanup on object destruction."""
        self.cleanup_instance()

    def _get_available_files(self, file_paths: List[Path]) -> List[Path]:
        """Filter files to only include those available for processing."""
        if not self.document_coordinator:
            return file_paths  # No coordination, process all files

        available_files = []
        for file_path in file_paths:
            if self.document_coordinator.can_process_file(file_path):
                available_files.append(file_path)

        return available_files

    def _process_single_file_with_coordination(self, file_path: Path) -> List[DocumentChunk]:
        """Process a single file with document coordination to prevent contention."""
        logger.debug(f"Processing file with coordination: {file_path}")

        # Try to acquire lock on the file
        if self.document_coordinator and not self.document_coordinator.acquire_file_lock(file_path):
            logger.info(f"File {file_path} is being processed by another instance, skipping")
            return []

        try:
            # Process the file
            chunks = self._process_single_file(file_path)

            # Mark file as completed
            if self.document_coordinator:
                self.document_coordinator.mark_file_completed(file_path)

            return chunks

        except Exception as e:
            # Mark file as failed and release lock
            if self.document_coordinator:
                self.document_coordinator.mark_file_failed(file_path, str(e))
            raise
        finally:
            # Always release the lock
            if self.document_coordinator:
                self.document_coordinator.release_file_lock(file_path)

    def _process_single_file(self, file_path: Path) -> List[DocumentChunk]:
        """Process a single file and return its chunks."""
        logger.debug(f"Processing file: {file_path}")

        # Extract text based on document type
        text = self._extract_text(file_path)

        # Split into chunks
        chunk_contents = self._split_text(text)

        # Create DocumentChunk objects
        chunks = []
        for i, content in enumerate(chunk_contents):
            chunk = DocumentChunk.create(
                content=content,
                source=str(file_path),
                metadata={
                    "type": self.config.doctype,
                    "chunk_index": i,
                    "chunking_strategy": self.config.chunking_strategy,
                    "file_size": file_path.stat().st_size if file_path.exists() else 0,
                },
            )
            chunks.append(chunk)

        return chunks

    def _extract_text(self, file_path: Path) -> str:
        """Extract text from a file based on its type."""
        if self.config.doctype == "json":
            with open(file_path, "r") as f:
                data = json.load(f)
            return str(data.get("text", json.dumps(data)))  # Ensure string return

        elif self.config.doctype == "pdf":
            text = ""
            with open(file_path, "rb") as file:
                reader = pypdf.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text()
            return text

        elif self.config.doctype == "txt":
            with open(file_path, "r", encoding="utf-8") as file:
                return file.read()

        elif self.config.doctype == "pptx":
            return self._extract_text_from_pptx(file_path)

        else:
            raise ValueError(f"Unsupported document type: {self.config.doctype}")

    def _extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint file."""
        prs = Presentation(str(file_path))  # Convert Path to str
        text_parts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
                elif hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text_parts.append(cell.text)

        return "\n".join(text_parts)

    def _split_text(self, text: str) -> List[str]:
        """Split text into chunks based on the configured strategy."""
        if self.config.chunking_strategy == "semantic":
            return self._semantic_chunking(text)
        elif self.config.chunking_strategy == "fixed":
            return self._fixed_chunking(text)
        elif self.config.chunking_strategy == "sentence":
            return self._sentence_chunking(text)
        else:
            raise ValueError(f"Unknown chunking strategy: {self.config.chunking_strategy}")

    def _semantic_chunking(self, text: str) -> List[str]:
        """Perform semantic chunking using embeddings."""
        # TODO: Implement semantic chunking with embeddings
        # For now, fall back to fixed chunking
        logger.warning("Semantic chunking not yet implemented, using fixed chunking")
        return self._fixed_chunking(text)

    def _fixed_chunking(self, text: str) -> List[str]:
        """Split text into fixed-size chunks."""
        chunk_size = self.config.chunk_size
        return [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)]

    def _sentence_chunking(self, text: str) -> List[str]:
        """Split text by sentences, respecting chunk size limits."""
        sentences = re.split(r"(?<=[.!?]) +", text)
        chunks = []
        current_chunk = ""

        for sentence in sentences:
            if len(current_chunk) + len(sentence) > self.config.chunk_size:
                if current_chunk:
                    chunks.append(current_chunk)
                current_chunk = sentence
            else:
                current_chunk += (" " if current_chunk else "") + sentence

        if current_chunk:
            chunks.append(current_chunk)

        return chunks

    def get_stats(self) -> EmbeddingStats:
        """Get processing statistics."""
        return self.stats
