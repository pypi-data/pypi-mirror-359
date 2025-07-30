from typing import Union

from loguru import logger
from pptx import Presentation

from datamax.parser.base import BaseLife, MarkdownOutputVo
from datamax.utils.lifecycle_types import LifeType


class PptxParser(BaseLife):
    def __init__(self, file_path: Union[str, list],domain: str = "Technology"):
        super().__init__(domain = domain)
        self.file_path = file_path

    @staticmethod
    def read_ppt_file(file_path: str):
        try:
            content = ""
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        content += shape.text + "\n"
            return content
        except Exception:
            raise

    def parse(self, file_path: str) -> MarkdownOutputVo:
        # —— 生命周期：开始处理 PPTX —— #
        lc_start = self.generate_lifecycle(
            source_file=file_path,
            domain=self.domain,
            usage_purpose="Documentation",
            life_type=LifeType.DATA_PROCESSING,
        )
        logger.debug("⚙️ DATA_PROCESSING 生命周期已生成")

        try:
            extension = self.get_file_extension(file_path)
            content = self.read_ppt_file(file_path=file_path)
            mk_content = content

            # —— 生命周期：处理完成 —— #
            lc_end = self.generate_lifecycle(
                source_file=file_path,
                domain=self.domain,
                usage_purpose="Documentation",
                life_type=LifeType.DATA_PROCESSED,
            )
            logger.debug("⚙️ DATA_PROCESSED 生命周期已生成")

            output_vo = MarkdownOutputVo(extension, mk_content)
            output_vo.add_lifecycle(lc_start)
            output_vo.add_lifecycle(lc_end)
            return output_vo.to_dict()

        except Exception as e:
            # —— 生命周期：处理失败 —— #
            lc_fail = self.generate_lifecycle(
                source_file=file_path,
                domain=self.domain,
                usage_purpose="Documentation",
                life_type=LifeType.DATA_PROCESS_FAILED,
            )
            logger.debug("⚙️ DATA_PROCESS_FAILED 生命周期已生成")

            raise Exception(
                {
                    "error": str(e),
                    "file_path": file_path,
                    "lifecycle": [lc_fail.to_dict()],
                }
            )
