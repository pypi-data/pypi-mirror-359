# ======================================================================
#  extractors.py  –  Unified text/metadata extractors
# ======================================================================
#
#  ### Iterator cheat-sheet
#  | Iterator              | What it yields                                                      |
#  | --------------------- | ------------------------------------------------------------------- |
#  | **git_repo_iter**     | Every text/config/code file in a GitHub / GitLab repo               |
#  | **confluence_iter**   | Every page in one or more Confluence spaces                         |
#  | **website_iter**      | All pages in a website (BFS within same domain, optional depth)     |
#  | **local_iter**        | Markdown from PDFs, DOCX, PPTX, XLSX, TXT, HTML, CSV, images,       |
#  |                     | audio, ZIPs, EPUB, JSON, XML, and more (recursive, via MarkItDown)   |
#
#  Power by MarkItDown:
#    - PDF, Word, PowerPoint, Excel, CSV, JSON, XML
#    - Images (EXIF + OCR), Audio (EXIF + transcription)
#    - HTML, Markdown, plain text, ZIPs (all inner contents), EPUB
#    - YouTube URLs (with transcript/extraction)
#
#  Optional installs
#  -----------------
#  # pip install markitdown
#  # pip install gitpython requests pandas tqdm
#  # pip install atlassian-python-api beautifulsoup4 lxml
#  # pip install pdfplumber python-docx python-pptx
#
#  Env vars for secrets (examples):
#     GITHUB_TOKEN, GITLAB_TOKEN, CONFLUENCE_TOKEN, CONF_USER, CONF_URL
# ======================================================================


from __future__ import annotations
from dataclasses import dataclass, asdict
from typing import Iterable, Optional, Set, Union, List, Dict, Any
from pathlib import Path
from datetime import datetime, UTC, timezone
from time import strftime
from collections import deque
from urllib.parse import urljoin, urlparse
import re, tempfile, os, shutil, requests, pandas as pd
import tempfile, shutil, requests, os
from git import Repo, InvalidGitRepositoryError, GitCommandError
from atlassian import Confluence
from bs4 import BeautifulSoup
from tqdm.auto import tqdm
import pdfplumber, docx
from markitdown import MarkItDown
from pptx import Presentation
import re


# ──────────────────────────  common record  ──────────────────────────
@dataclass
class DocumentRecord:
    source: str
    filename: str
    title: str          | None
    description: str    | None
    author: str         | None
    content: str
    path_or_id: str     | None
    url: str            | None
    last_modified: str  | None
    size_bytes: int     | None
    def to_dict(self): return asdict(self)


# ════════════════════════════  GIT REPOS  ════════════════════════════

# ───────────────────── helper – robust “clean & clone” ─────────────────────
def _prepare_repo_dir(clone_url: str, repo_dir: Path) -> Repo:
    """
    Ensure *repo_dir* is a fresh, valid Git repo.

    • If it's already a healthy repo → reuse it.
    • If it's corrupted → move it aside and reclone.
    • Windows-safe: handles file-locks and existing “_old” dirs.
    """
    if repo_dir.exists():
        try:
            return Repo(repo_dir) 
        except InvalidGitRepositoryError:
            print(f"[info] {repo_dir} is broken; refreshing…")

            # Build a unique “…_old” path
            alt = repo_dir.with_name(repo_dir.name + "_old")
            if alt.exists():
                # If stale, nuke it first
                shutil.rmtree(alt, ignore_errors=True)

            try:
                repo_dir.rename(alt)
            except (PermissionError, FileExistsError):
                # Last-resort uniqueness
                ts_alt = repo_dir.with_name(
                    f"{repo_dir.name}_old_{datetime.now(UTC).strftime('%Y%m%d%H%M%S')}"
                )
                try:
                    repo_dir.rename(ts_alt)
                except Exception as e:
                    print(f"[warn] rename failed ({e}); force-deleting")
                    shutil.rmtree(repo_dir, ignore_errors=True)

    # Now the path is definitely free
    return Repo.clone_from(clone_url, repo_dir, depth=1, single_branch=True)


# ───────────────────────────────  main iterator  ───────────────────────────
def git_repo_iter(
    repo_url: str,
    token: Optional[str] = None,
    *,
    clone_dir: Union[str, Path] = Path(tempfile.gettempdir()) / "git-mirror",
    include_ext: Optional[Set[str]] = None,
    delete_after: bool = True,
) -> Iterable[DocumentRecord]:
    """
    Stream `DocumentRecord` for every text/config file in a GitHub or GitLab repo.

    Parameters
    ----------
    repo_url : str
        HTTPS clone URL (`…github.com/owner/repo.git` or `…gitlab.com/group/repo.git`)
    token : str | None
        Personal-access token; omit for public repos.
    clone_dir : Path
        Parent folder under which a temp checkout is created.
    include_ext : set[str] | None
        File extensions to include (defaults to common text/code types).
    delete_after : bool
        Remove the temp repo folder once iteration completes.
    """
    include_ext = include_ext or {
        ".md", ".markdown", ".txt",
        ".py", ".ipynb", ".r", ".R", ".sh",
        ".yaml", ".yml", ".json", ".toml", ".ini", ".cfg", ".conf",
        ".csv", ".tsv", ".sql",
        ".html", ".htm", ".xml", ".rst",
    }

    source = (
        "github"  if "github.com"  in repo_url else
        "gitlab"  if "gitlab.com"  in repo_url else
        "git"
    )

    clone_dir = Path(clone_dir).expanduser()
    repo_name = repo_url.rstrip("/").split("/")[-1].replace(".git", "")
    repo_dir  = clone_dir / repo_name

    # Build auth-injected URL when a token is supplied
    clone_url = repo_url
    if token and repo_url.startswith("https://"):
        if "github.com" in repo_url:
            clone_url = repo_url.replace("https://", f"https://{token}@")
        elif "gitlab.com" in repo_url:
            clone_url = repo_url.replace("https://", f"https://oauth2:{token}@")

    # Clone or refresh corrupted checkouts
    try:
        repo = _prepare_repo_dir(clone_url, repo_dir)
    except GitCommandError as err:
        print(f"[error] git clone failed for {repo_url}: {err}")
        return  # abort iteration gracefully

    repo_descr = _repo_description(repo_url, token)

    try:
        for path in repo_dir.rglob("*"):
            if not (path.is_file() and path.suffix.lower() in include_ext):
                continue
            try:
                text = path.read_text(encoding="utf-8", errors="replace")
            except Exception as exc:
                print(f"[warn] skip {path}: {exc}")
                continue

            author = _file_author(repo, path)
            base   = repo_url.replace(".git", "")
            view   = (
                f"{base}/-/blob/main/{path.relative_to(repo_dir)}"
                if "gitlab.com" in repo_url
                else f"{base}/blob/main/{path.relative_to(repo_dir)}"
            )

            yield DocumentRecord(
                source=source,
                filename=path.name,
                title=path.stem,
                description=repo_descr,
                author=author,
                content=text,
                path_or_id=str(path.relative_to(repo_dir)),
                url=view,
                last_modified=datetime.fromtimestamp(
                    path.stat().st_mtime, tz=UTC
                ).isoformat(),
                size_bytes=path.stat().st_size,
            )
    finally:
        if delete_after and repo_dir.exists():
            shutil.rmtree(repo_dir, ignore_errors=True)


# ────────────────────────────  helper functions  ──────────────────────────
def _repo_description(repo_url: str, token: Optional[str]) -> Optional[str]:
    try:
        hdr = {"Authorization": f"Bearer {token}"} if token else {}
        if "github.com" in repo_url:
            owner_repo = repo_url.replace("https://github.com/", "").replace(".git", "")
            r = requests.get(f"https://api.github.com/repos/{owner_repo}", headers=hdr, timeout=15)
            return r.json().get("description") if r.ok else None
        if "gitlab.com" in repo_url:
            pth = repo_url.replace("https://gitlab.com/", "").replace(".git", "")
            api = f"https://gitlab.com/api/v4/projects/{requests.utils.quote(pth, safe='')}"
            r = requests.get(api, headers=hdr, timeout=15)
            return r.json().get("description") if r.ok else None
    except Exception:
        pass
    return None


def _file_author(repo: Repo, fp: Path) -> Optional[str]:
    try:
        rel = str(fp.relative_to(repo.working_tree_dir))
        commit = next(repo.iter_commits(paths=rel, max_count=1), None)
        return commit.author.name if commit else None
    except Exception:
        return None


# ════════════════════════════  CONFLUENCE  ════════════════════════════
def confluence_iter(
    base_url: str,
    username: str,
    api_token: str,
    space_keys: List[str],
    *,
    limit: int | None = None,
    plain_text: bool = False,
) -> Iterable[DocumentRecord]:
    """
    Stream pages from Confluence Cloud / DC.
    """
    cf = Confluence(url=base_url, username=username, password=api_token)

    for space in space_keys:
        s_meta = cf.get_space(space, expand="description.plain")
        s_descr = (s_meta.get("description", {}).get("plain", {}).get("value"))

        start, fetched, size = 0, 0, 100
        pbar = tqdm(desc=f"[{space}] pages", unit="page")

        while True:
            pages = cf.get_all_pages_from_space(space, start=start, limit=size,
                                                expand="version,body.storage")
            if not pages:
                break
            for p in pages:
                fetched += 1; pbar.update(1)
                if limit and fetched > limit:
                    break

                pid, html = p["id"], p["body"]["storage"]["value"]
                cnt = BeautifulSoup(html, "html.parser").get_text("\n") if plain_text else html
                yield DocumentRecord(
                    source="confluence",
                    filename=f"{space}-{pid}",
                    title=p["title"],
                    description=s_descr,
                    author=p["version"]["by"]["displayName"],
                    content=cnt,
                    path_or_id=pid,
                    url=f"{base_url.rstrip('/')}/pages/viewpage.action?pageId={pid}",
                    last_modified=datetime.fromisoformat(p["version"]["when"]).astimezone(UTC).isoformat(),
                    size_bytes=len(html.encode()),
                )
            if limit and fetched >= limit:
                break
            start += size
        pbar.close()


# ════════════════════════════  WEBSITE CRAWL  ═════════════════════════
# ───────────────────── helper – read <meta> tags ─────────────────────
def _meta(soup: BeautifulSoup, key: str) -> str | None:
    """
    Return the content of the first <meta name="<key>" …> or
    <meta property="og:<key>" …> tag, if present.
    """
    tag = (
        soup.find("meta", attrs={"name": key})
        or soup.find("meta", attrs={"property": f"og:{key}"})
    )
    return tag.get("content", "").strip() if tag and tag.has_attr("content") else None

def website_iter(
    roots: List[str],
    *,
    crawl_depth: int | None = None,
) -> Iterable[DocumentRecord]:
    """
    Breadth-first crawl within each domain and yield one `DocumentRecord`
    per page.
    """
    def _clean(txt: str) -> str:
        txt = re.sub(r"\s+\n", "\n", txt)
        return re.sub(r"\n{2,}", "\n", txt).strip()

    seen: set[str] = set()
    for root in roots:
        q: deque[tuple[str,int]] = deque([(root,0)])
        host = urlparse(root).netloc

        while q:
            url, d = q.popleft()
            if url in seen: continue
            seen.add(url)

            try:
                html = requests.get(url, timeout=15).text
            except Exception:
                continue
            soup = BeautifulSoup(html, "lxml")

            yield DocumentRecord(
                source=host,
                filename=url,
                title=soup.title.string.strip() if soup.title else None,
                description=_meta(soup, "description"),
                author=_meta(soup, "author"),
                content=_clean(soup.get_text(" ", strip=True)),
                path_or_id=None,
                url=url,
                last_modified=None,
                size_bytes=len(html.encode()),
            )

            if crawl_depth is None or d < crawl_depth:
                for a in soup.find_all("a", href=True):
                    nxt = urljoin(url, a["href"].split("#")[0])
                    if urlparse(nxt).netloc == host and nxt not in seen:
                        q.append((nxt, d+1))


# ════════════════════════════  LOCAL FILES  ═══════════════════════════
def local_iter(
    paths: list[str | Path],
    *,
    include_ext: Optional[Set[str]] = None,
    # new optional args for MarkItDown:
    llm_client: Any = None,     # e.g. OpenAI(api_key=…)
    llm_model: str | None = None,
) -> Iterable[DocumentRecord]:
    """
    Recursively crawl `paths` and yield a DocumentRecord for each file.
    Uses MarkItDown to convert supported types to Markdown via an LLM,
    falling back to per-extension extractors if that fails.
    """
    # 1) Build your MarkItDown instance once
    md = MarkItDown(
        llm_client=llm_client,
        llm_model=llm_model or "gpt-4"
    )

    # 2) Default extension filter (optional–you can override)
    include_ext = include_ext or {
        ".md", ".markdown", ".txt",
        ".py", ".ipynb", ".r", ".R", ".sh",
        ".yaml", ".yml", ".json", ".toml", ".ini", ".cfg", ".conf",
        ".csv", ".tsv", ".sql",
        ".html", ".htm", ".xml", ".rst",
        ".pdf", ".docx", ".pptx", ".xlsx", ".xls",
        ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff",
        ".mp3", ".wav", ".ogg",
        ".zip", ".epub"
    }

    for p in paths:
        p = Path(p)
        if p.is_dir():
            for f in p.rglob("*"):
                if f.suffix.lower() in include_ext:
                    recs = _read_local(f, md)
                    if recs:
                        # handle multi‐record (e.g. ZIP) vs single
                        if isinstance(recs, Iterable) and not isinstance(recs, (str, DocumentRecord)):
                            for rec in recs:
                                yield rec
                        else:
                            yield recs

        elif p.exists() and p.suffix.lower() in include_ext:
            recs = _read_local(p, md)
            if recs:
                if isinstance(recs, Iterable) and not isinstance(recs, (str, DocumentRecord)):
                    for rec in recs:
                        yield rec
                else:
                    yield recs

        else:
            print(f"[warn] unsupported or missing: {p}")
            

# --------------- local readers ----------------
def _clean(text: str) -> str:
    """
    Collapse super-long whitespace / blank lines.
    """
    text = re.sub(r"\s+\n", "\n", text)   # strip trailing spaces on each line
    return re.sub(r"\n{2,}", "\n", text).strip()

def _read_local(
    path: Path,
    md: MarkItDown
) -> DocumentRecord | Iterable[DocumentRecord] | None:
    """
    1) Try md.convert(...) to produce markdown via LLM
    2) On failure, fall back to per-filetype extraction
    """
    ext = path.suffix.lower()

    # --- 1) Primary: MarkItDown via LLM ---
    try:
        result = md.convert(str(path))
        # result.text_content holds the Markdown string
        md_text = getattr(result, "text_content", result)

        # ZIP/Archive support? If result is dict-like, combine:
        if isinstance(md_text, dict):
            combined = "\n\n".join(
                f"### {fname}\n\n{content}"
                for fname, content in md_text.items()
            )
            return DocumentRecord(
                source="local",
                filename=path.name,
                title=None, description=None, author=None,
                content=combined,
                path_or_id=str(path),
                url=None,
                last_modified=None,
                size_bytes=path.stat().st_size
            )

        return DocumentRecord(
            source="local",
            filename=path.name,
            title=None, description=None, author=None,
            content=md_text,
            path_or_id=str(path),
            url=None,
            last_modified=None,
            size_bytes=path.stat().st_size
        )

    except Exception as e:
        print(f"[warn] MarkItDown failed for {path}: {e}")

    # --- 2) Fallback: per‐extension extractor ---
    try:
        if ext == ".pdf":
            with pdfplumber.open(path) as pdf:
                meta = pdf.metadata or {}
                txt  = "\n".join(p.extract_text() or "" for p in pdf.pages)
            return DocumentRecord(
                "local", path.name, meta.get("Title"), None, meta.get("Author"),
                _clean(txt), str(path), None,
                _pdf_date(meta.get("CreationDate")), path.stat().st_size
            )

        if ext == ".docx":
            doc  = docx.Document(path)
            core = doc.core_properties
            txt  = "\n".join(p.text for p in doc.paragraphs)
            return DocumentRecord(
                "local", path.name, core.title, None, core.author,
                _clean(txt), str(path), None,
                _py_dt(core.created), path.stat().st_size
            )

        if ext == ".pptx":
            prs   = Presentation(path)
            core  = prs.core_properties
            slides_txt = [
                shape.text for sl in prs.slides for shape in sl.shapes
                if hasattr(shape, "text")
            ]
            return DocumentRecord(
                "local", path.name, core.title, None, core.author,
                _clean("\n".join(slides_txt)), str(path), None,
                _py_dt(core.created), path.stat().st_size
            )

        if ext in {".html", ".htm"}:
            soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "lxml")
            return DocumentRecord(
                "local", path.name,
                soup.title.string.strip() if soup.title else None,
                None, _meta(soup, "author"),
                _clean(soup.get_text(" ", strip=True)), str(path), None,
                None, path.stat().st_size
            )

        if ext in {
            ".txt", ".md", ".markdown", ".rst", ".py", ".sh", ".r", ".R",
            ".yaml", ".yml", ".json", ".toml", ".ini", ".cfg", ".conf",
            ".csv", ".tsv", ".sql", ".xml"
        }:
            txt = path.read_text(encoding="utf-8", errors="ignore")
            return DocumentRecord(
                "local", path.name, None, None, None,
                _clean(txt), str(path), None,
                None, path.stat().st_size
            )

    except Exception as ex:
        print(f"[warn] fallback extractor failed for {path}: {ex}")

    return None
