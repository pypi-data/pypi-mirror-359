# ======================================================================
#  extractors.py  –  Unified text/metadata extractors
# ======================================================================
#
#  ### Iterator cheat-sheet
#  | Iterator              | What it yields                                                      |
#  | --------------------- | ------------------------------------------------------------------- |
#  | **git_repo_iter**     | Every text/config/code file in a GitHub / GitLab repo               |
#  | **confluence_iter**   | Every page in one or more Confluence spaces                         |
#  | **website_iter**      | All pages in a website (BFS within same domain, optional depth)     |
#  | **local_iter**        | Markdown from PDFs, DOCX, PPTX, XLSX, TXT, HTML, CSV, images,       |
#  |                     | audio, ZIPs, EPUB, JSON, XML, and more (recursive, via MarkItDown)   |
#
#  Power by MarkItDown:
#    - PDF, Word, PowerPoint, Excel, CSV, JSON, XML
#    - Images (EXIF + OCR), Audio (EXIF + transcription)
#    - HTML, Markdown, plain text, ZIPs (all inner contents), EPUB
#    - YouTube URLs (with transcript/extraction)
#
#  Optional installs
#  -----------------
#  # pip install markitdown
#  # pip install gitpython requests pandas tqdm
#  # pip install atlassian-python-api beautifulsoup4 lxml
#  # pip install pdfplumber python-docx python-pptx
#
#  Env vars for secrets (examples):
#     GITHUB_TOKEN, GITLAB_TOKEN, CONFLUENCE_TOKEN, CONF_USER, CONF_URL
# ======================================================================


from __future__ import annotations

# ──────────────────────── standard library ────────────────────────────
from collections import deque
from dataclasses import dataclass, asdict
from datetime import UTC, datetime, timedelta, timezone
from pathlib import Path
from typing import (
    Any,
    Iterable,
    List,
    Optional,
    Set,
    Tuple,
    Union,
)
from urllib.parse import urljoin, urlparse, urlunparse
import os
import re
import shutil
import tempfile
import textwrap
import unicodedata

# ───────────────────────── third-party libs ───────────────────────────
import requests
import pandas as pd
from atlassian import Confluence
from bs4 import BeautifulSoup
from git import GitCommandError, InvalidGitRepositoryError, Repo
from markitdown import MarkItDown
from pptx import Presentation
from tqdm.auto import tqdm
import pdfplumber
import docx




# ──────────────────────────  common record  ──────────────────────────
@dataclass
class DocumentRecord:
    source: str
    filename: str
    title: str          | None
    description: str    | None
    author: str         | None
    content: str
    path_or_id: str     | None
    url: str            | None
    last_modified: str  | None
    size_bytes: int     | None
    def to_dict(self): return asdict(self)


# ════════════════════════════  GIT REPOS  ════════════════════════════

# ───────────────────── helper – robust “clean & clone” ─────────────────────
def _prepare_repo_dir(clone_url: str, repo_dir: Path) -> Repo:
    """
    Ensure *repo_dir* is a fresh, valid Git repo.

    • If it's already a healthy repo → reuse it.
    • If it's corrupted → move it aside and reclone.
    • Windows-safe: handles file-locks and existing “_old” dirs.
    """
    if repo_dir.exists():
        try:
            return Repo(repo_dir) 
        except InvalidGitRepositoryError:
            print(f"[info] {repo_dir} is broken; refreshing…")

            # Build a unique “…_old” path
            alt = repo_dir.with_name(repo_dir.name + "_old")
            if alt.exists():
                # If stale, nuke it first
                shutil.rmtree(alt, ignore_errors=True)

            try:
                repo_dir.rename(alt)
            except (PermissionError, FileExistsError):
                # Last-resort uniqueness
                ts_alt = repo_dir.with_name(
                    f"{repo_dir.name}_old_{datetime.now(UTC).strftime('%Y%m%d%H%M%S')}"
                )
                try:
                    repo_dir.rename(ts_alt)
                except Exception as e:
                    print(f"[warn] rename failed ({e}); force-deleting")
                    shutil.rmtree(repo_dir, ignore_errors=True)

    # Now the path is definitely free
    return Repo.clone_from(clone_url, repo_dir, depth=1, single_branch=True)


# ───────────────────────────────  main iterator  ───────────────────────────
def git_repo_iter(
    repo_url: str,
    token: Optional[str] = None,
    *,
    clone_dir: Union[str, Path] = Path(tempfile.gettempdir()) / "git-mirror",
    include_ext: Optional[Set[str]] = None,
    delete_after: bool = True,
) -> Iterable[DocumentRecord]:
    """
    Stream `DocumentRecord` for every text/config file in a GitHub or GitLab repo.

    Parameters
    ----------
    repo_url : str
        HTTPS clone URL (`…github.com/owner/repo.git` or `…gitlab.com/group/repo.git`)
    token : str | None
        Personal-access token; omit for public repos.
    clone_dir : Path
        Parent folder under which a temp checkout is created.
    include_ext : set[str] | None
        File extensions to include (defaults to common text/code types).
    delete_after : bool
        Remove the temp repo folder once iteration completes.
    """
    include_ext = include_ext or {
        ".md", ".markdown", ".txt",
        ".py", ".ipynb", ".r", ".R", ".sh",
        ".yaml", ".yml", ".json", ".toml", ".ini", ".cfg", ".conf",
        ".csv", ".tsv", ".sql",
        ".html", ".htm", ".xml", ".rst",
    }

    source = (
        "github"  if "github.com"  in repo_url else
        "gitlab"  if "gitlab.com"  in repo_url else
        "git"
    )

    clone_dir = Path(clone_dir).expanduser()
    repo_name = repo_url.rstrip("/").split("/")[-1].replace(".git", "")
    repo_dir  = clone_dir / repo_name

    # Build auth-injected URL when a token is supplied
    clone_url = repo_url
    if token and repo_url.startswith("https://"):
        if "github.com" in repo_url:
            clone_url = repo_url.replace("https://", f"https://{token}@")
        elif "gitlab.com" in repo_url:
            clone_url = repo_url.replace("https://", f"https://oauth2:{token}@")

    # Clone or refresh corrupted checkouts
    try:
        repo = _prepare_repo_dir(clone_url, repo_dir)
    except GitCommandError as err:
        print(f"[error] git clone failed for {repo_url}: {err}")
        return  # abort iteration gracefully

    repo_descr = _repo_description(repo_url, token)

    try:
        for path in repo_dir.rglob("*"):
            if not (path.is_file() and path.suffix.lower() in include_ext):
                continue
            try:
                text = path.read_text(encoding="utf-8", errors="replace")
            except Exception as exc:
                print(f"[warn] skip {path}: {exc}")
                continue

            author = _file_author(repo, path)
            base   = repo_url.replace(".git", "")
            view   = (
                f"{base}/-/blob/main/{path.relative_to(repo_dir)}"
                if "gitlab.com" in repo_url
                else f"{base}/blob/main/{path.relative_to(repo_dir)}"
            )

            yield DocumentRecord(
                source=source,
                filename=path.name,
                title=path.stem,
                description=repo_descr,
                author=author,
                content=text,
                path_or_id=str(path.relative_to(repo_dir)),
                url=view,
                last_modified=datetime.fromtimestamp(
                    path.stat().st_mtime, tz=UTC
                ).isoformat(),
                size_bytes=path.stat().st_size,
            )
    finally:
        if delete_after and repo_dir.exists():
            shutil.rmtree(repo_dir, ignore_errors=True)


# ────────────────────────────  helper functions  ──────────────────────────
def _repo_description(repo_url: str, token: Optional[str]) -> Optional[str]:
    try:
        hdr = {"Authorization": f"Bearer {token}"} if token else {}
        if "github.com" in repo_url:
            owner_repo = repo_url.replace("https://github.com/", "").replace(".git", "")
            r = requests.get(f"https://api.github.com/repos/{owner_repo}", headers=hdr, timeout=15)
            return r.json().get("description") if r.ok else None
        if "gitlab.com" in repo_url:
            pth = repo_url.replace("https://gitlab.com/", "").replace(".git", "")
            api = f"https://gitlab.com/api/v4/projects/{requests.utils.quote(pth, safe='')}"
            r = requests.get(api, headers=hdr, timeout=15)
            return r.json().get("description") if r.ok else None
    except Exception:
        pass
    return None


def _file_author(repo: Repo, fp: Path) -> Optional[str]:
    try:
        rel = str(fp.relative_to(repo.working_tree_dir))
        commit = next(repo.iter_commits(paths=rel, max_count=1), None)
        return commit.author.name if commit else None
    except Exception:
        return None


# ════════════════════════════  CONFLUENCE  ════════════════════════════
def confluence_iter(
    base_url: str,
    username: str,
    api_token: str,
    space_keys: List[str],
    *,
    limit: int | None = None,
    plain_text: bool = False,
) -> Iterable[DocumentRecord]:
    """
    Stream pages from Confluence Cloud / DC.
    """
    cf = Confluence(url=base_url, username=username, password=api_token)

    for space in space_keys:
        s_meta = cf.get_space(space, expand="description.plain")
        s_descr = (s_meta.get("description", {}).get("plain", {}).get("value"))

        start, fetched, size = 0, 0, 100
        pbar = tqdm(desc=f"[{space}] pages", unit="page")

        while True:
            pages = cf.get_all_pages_from_space(space, start=start, limit=size,
                                                expand="version,body.storage")
            if not pages:
                break
            for p in pages:
                fetched += 1; pbar.update(1)
                if limit and fetched > limit:
                    break

                pid, html = p["id"], p["body"]["storage"]["value"]
                cnt = BeautifulSoup(html, "html.parser").get_text("\n") if plain_text else html
                yield DocumentRecord(
                    source="confluence",
                    filename=f"{space}-{pid}",
                    title=p["title"],
                    description=s_descr,
                    author=p["version"]["by"]["displayName"],
                    content=cnt,
                    path_or_id=pid,
                    url=f"{base_url.rstrip('/')}/pages/viewpage.action?pageId={pid}",
                    last_modified=datetime.fromisoformat(p["version"]["when"]).astimezone(UTC).isoformat(),
                    size_bytes=len(html.encode()),
                )
            if limit and fetched >= limit:
                break
            start += size
        pbar.close()


# ════════════════════════════  WEBSITE CRAWL  ═════════════════════════
# ──────────────────── stdlib / third-party imports ────────────────────
# ──────────────────────── helper – read <meta> tags ───────────────────
def _meta(soup: BeautifulSoup, key: str) -> Optional[str]:
    """
    Return the content of the first matching meta tag:
        <meta name="<key>" …>           or
        <meta property="og:<key>" …>    or
        <meta property="twitter:<key>" …>
    """
    tag = (
        soup.find("meta", attrs={"name": key})
        or soup.find("meta", attrs={"property": f"og:{key}"})
        or soup.find("meta", attrs={"property": f"twitter:{key}"})
    )
    return tag.get("content", "").strip() if tag and tag.has_attr("content") else None


# ──────────────────────── helper – canonicalise URL ────────────────────
def _canonical(url: str) -> str:
    """
    Normalise URL so duplicates (`/` vs ``) collapse to one:

      • strips leading 'www.'
      • removes trailing slash (except `/` root)
      • drops URL fragment  (#section)
      • keeps scheme so http/https stay distinct
    """
    u = urlparse(url)
    host = u.netloc.lstrip("www.")
    path = (u.path or "/").rstrip("/") or "/"
    return urlunparse((u.scheme, host, path, "", "", ""))


# ─────────────────────── helper – standard text clean ──────────────────
def _clean(
    text: str,
    *,
    normalize_unicode: bool = True,
    collapse_whitespace: bool = True,
    strip_line_ends: bool = True,
    dedent: bool = True,
) -> str:
    """Unicode-safe, whitespace-aware cleaner (see docstring in previous reply)."""
    if normalize_unicode:
        text = unicodedata.normalize("NFKC", text)
        text = text.replace("\u00A0", " ")                       # NBSP
        text = re.sub(r"[\u200B-\u200D\u2060\uFEFF]", "", text)  # zero-width

    text = text.replace("\r\n", "\n").replace("\r", "\n")

    if strip_line_ends:
        text = "\n".join(line.rstrip() for line in text.splitlines())

    if dedent:
        text = textwrap.dedent(text)

    if collapse_whitespace:
        text = re.sub(r"[ \t]+", " ", text)      # multi-space → 1
        text = re.sub(r"\n{3,}", "\n\n", text)   # >2 blank lines → 1

    return text.strip()


# ═══════════════════════════  WEBSITE CRAWLER  ═════════════════════════
def website_iter(
    roots: List[str],
    *,
    crawl_depth: Optional[int] = None,
    user_agent: str = "Mozilla/5.0 (compatible; extractor/1.0)",
) -> Iterable["DocumentRecord"]:
    """
    Breadth-first crawl within each domain and yield a `DocumentRecord`
    per HTML page.

    • **Deduplicates** by canonical URL (so / and /index.html don’t duplicate)
    • **Stops** at *crawl_depth* relative to each root (None = unlimited)
    • **Skips** non-HTML responses (content-type gate)
    """
    seen: set[str] = set()
    hdrs = {"User-Agent": user_agent}

    for root in roots:
        queue: deque[Tuple[str, int]] = deque([(root, 0)])
        host = urlparse(root).netloc.lstrip("www.")

        while queue:
            url, depth = queue.popleft()
            canon = _canonical(url)
            if canon in seen:
                continue
            seen.add(canon)

            try:
                resp = requests.get(url, headers=hdrs, timeout=15)
                ctype = resp.headers.get("content-type", "")
                if "html" not in ctype:
                    continue
                html = resp.text
            except Exception:
                continue

            soup = BeautifulSoup(html, "lxml")

            yield DocumentRecord(
                source=host,
                filename=canon,
                title=soup.title.string.strip() if soup.title else None,
                description=_meta(soup, "description"),
                author=_meta(soup, "author"),
                content=_clean(soup.get_text(" ", strip=True)),
                path_or_id=None,
                url=canon,
                last_modified=datetime.now(UTC).isoformat(),  # or parse <meta>
                size_bytes=len(html.encode()),
            )

            # enqueue same-domain links
            if crawl_depth is None or depth < crawl_depth:
                for a in soup.find_all("a", href=True):
                    nxt = urljoin(url, a["href"].split("#")[0])
                    if urlparse(nxt).netloc.lstrip("www.") == host:
                        queue.append((nxt, depth + 1))


# ════════════════════════════  LOCAL FILES  ═══════════════════════════
def local_iter(
    paths: list[str | Path],
    *,
    include_ext: Optional[Set[str]] = None,
    # new optional args for MarkItDown:
    llm_client: Any = None,     # e.g. OpenAI(api_key=…)
    llm_model: str | None = None,
) -> Iterable[DocumentRecord]:
    """
    Recursively crawl `paths` and yield a DocumentRecord for each file.
    Uses MarkItDown to convert supported types to Markdown via an LLM,
    falling back to per-extension extractors if that fails.
    """
    # 1) Build your MarkItDown instance once
    md = MarkItDown(
        llm_client=llm_client,
        llm_model=llm_model or "gpt-4"
    )

    # 2) Default extension filter (optional–you can override)
    include_ext = include_ext or {
        ".md", ".markdown", ".txt",
        ".py", ".ipynb", ".r", ".R", ".sh",
        ".yaml", ".yml", ".json", ".toml", ".ini", ".cfg", ".conf",
        ".csv", ".tsv", ".sql",
        ".html", ".htm", ".xml", ".rst",
        ".pdf", ".docx", ".pptx", ".xlsx", ".xls",
        ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff",
        ".mp3", ".wav", ".ogg",
        ".zip", ".epub"
    }

    for p in paths:
        p = Path(p)
        if p.is_dir():
            for f in p.rglob("*"):
                if f.suffix.lower() in include_ext:
                    recs = _read_local(f, md)
                    if recs:
                        # handle multi‐record (e.g. ZIP) vs single
                        if isinstance(recs, Iterable) and not isinstance(recs, (str, DocumentRecord)):
                            for rec in recs:
                                yield rec
                        else:
                            yield recs

        elif p.exists() and p.suffix.lower() in include_ext:
            recs = _read_local(p, md)
            if recs:
                if isinstance(recs, Iterable) and not isinstance(recs, (str, DocumentRecord)):
                    for rec in recs:
                        yield rec
                else:
                    yield recs

        else:
            print(f"[warn] unsupported or missing: {p}")
            

# --------------- local readers ----------------
# ----------------------------------------------------------------------
# Helper – normalise python-docx / pptx core-property datetimes
# ----------------------------------------------------------------------
_ISO_RE = re.compile(
    r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}(?:\.\d+)?(?:Z|[+\-]\d{2}:\d{2})?$"
)

def _py_dt(raw) -> str | None:
    """
    Convert a python-docx / pptx datetime (or ISO-8601 string) to
    an ISO-8601 UTC string.

    Accepts:
      • `datetime` objects (tz-aware or naive)
      • ISO-8601 strings (with/without timezone)
      • Anything else → returns `None`
    """
    if raw is None:
        return None

    # Already datetime -----------------------------------------------
    if isinstance(raw, datetime):
        dt = raw
        if dt.tzinfo is None:                 # make naive → UTC
            dt = dt.replace(tzinfo=timezone.utc)
        return dt.astimezone(timezone.utc).isoformat()

    # ISO string (python-docx sometimes stores str) ------------------
    if isinstance(raw, str) and _ISO_RE.match(raw):
        try:
            dt = datetime.fromisoformat(
                raw.replace("Z", "+00:00")     # handle trailing Z
            )
            if dt.tzinfo is None:
                dt = dt.replace(tzinfo=timezone.utc)
            return dt.astimezone(timezone.utc).isoformat()
        except ValueError:
            return None

    # Unrecognised ----------------------------------------------------
    return None

# ----------------------------------------------------------------------
# Helper – parse PDF metadata date strings  (e.g. "D:20230715144300-05'00'")
# ----------------------------------------------------------------------
_pdf_dt_re = re.compile(
    r"""^D:(\d{4})(\d{2})(\d{2})          # YYYY MM DD
        (?:(\d{2})(\d{2})(\d{2}))?        # optional HH mm SS
        (?:
            ([\+\-Z])(\d{2})'?(\d{2})'?   # timezone  (+HH'mm' | -HH'mm' | Z)
        )?$
    """,
    re.VERBOSE,
)

def _pdf_date(raw: str | None) -> str | None:
    """
    Convert PDF metadata dates (D:YYYYMMDDHHmmSS+HH'mm') → ISO-8601.

    Returns None if parsing fails.
    """
    if not raw:
        return None

    m = _pdf_dt_re.match(raw)
    if not m:
        return None

    (
        year, mon, day,
        hour, minute, sec,
        tz_sign, tz_h, tz_m
    ) = m.groups()

    # Default to midnight if no time part
    hour   = int(hour   or 0)
    minute = int(minute or 0)
    sec    = int(sec    or 0)

    # Time-zone handling
    if tz_sign in ("+", "-"):
        offset = int(tz_h) * 60 + int(tz_m)
        if tz_sign == "-":
            offset = -offset
        tz = timezone(timedelta(minutes=offset))
    else:                 # 'Z' or missing → UTC
        tz = timezone.utc

    try:
        dt = datetime(
            int(year), int(mon), int(day),
            hour, minute, sec,
            tzinfo=tz
        )
        return dt.astimezone(timezone.utc).isoformat()
    except ValueError:
        return None

def _clean(
    text: str,
    *,
    normalize_unicode: bool = True,
    collapse_whitespace: bool = True,
    strip_line_ends: bool = True,
    dedent: bool = True,
) -> str:
    """
    Standard-clean raw text extracted from PDFs, HTML, Office docs, etc.

    Steps
    -----
    1. **Unicode normalisation** (NFKC) & removal of zero-width / NBSP chars
    2. **New-line normalisation**  (CRLF ⇢ LF)
    3. **Line-end trimming**        (rstrip each line)
    4. **Block dedent**             (textwrap.dedent)
    5. **Whitespace collapse**      (⇢ single spaces, ≤2 consecutive LFs)

    Parameters
    ----------
    text : str
        Raw extracted text.
    normalize_unicode, collapse_whitespace, strip_line_ends, dedent : bool
        Toggle individual cleaning steps for special cases.

    Returns
    -------
    str
        Cleaned, normalised text ready for chunking / embedding.
    """
    # 1️⃣  Unicode / control-char normalisation
    if normalize_unicode:
        text = unicodedata.normalize("NFKC", text)
        text = text.replace("\u00A0", " ")                 # NBSP
        text = re.sub(r"[\u200B-\u200D\u2060\uFEFF]", "", text)  # zero-width

    # 2️⃣  Standardise new-lines (CRLF/R ⇢ LF)
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    # 3️⃣  Trim trailing spaces on each line
    if strip_line_ends:
        text = "\n".join(line.rstrip() for line in text.splitlines())

    # 4️⃣  Remove common leading indentation
    if dedent:
        text = textwrap.dedent(text)

    # 5️⃣  Collapse whitespace
    if collapse_whitespace:
        text = re.sub(r"[ \t]+", " ", text)       # multiple spaces → one
        text = re.sub(r"\n{3,}", "\n\n", text)    # >2 blank lines → 1

    return text.strip()


def _read_local(
    path: Path,
    md: MarkItDown
) -> DocumentRecord | Iterable[DocumentRecord] | None:
    """
    1) Try md.convert(...) to produce markdown via LLM
    2) On failure, fall back to per-filetype extraction
    """
    ext = path.suffix.lower()

    # --- 1) Primary: MarkItDown via LLM ---
    try:
        result = md.convert(str(path))
        # result.text_content holds the Markdown string
        md_text = getattr(result, "text_content", result)

        # ZIP/Archive support? If result is dict-like, combine:
        if isinstance(md_text, dict):
            combined = "\n\n".join(
                f"### {fname}\n\n{content}"
                for fname, content in md_text.items()
            )
            return DocumentRecord(
                source="local",
                filename=path.name,
                title=None, description=None, author=None,
                content=combined,
                path_or_id=str(path),
                url=None,
                last_modified=None,
                size_bytes=path.stat().st_size
            )

        return DocumentRecord(
            source="local",
            filename=path.name,
            title=None, description=None, author=None,
            content=md_text,
            path_or_id=str(path),
            url=None,
            last_modified=None,
            size_bytes=path.stat().st_size
        )

    except Exception as e:
        print(f"[warn] MarkItDown failed for {path}: {e}")

    # --- 2) Fallback: per‐extension extractor ---
    try:
        if ext == ".pdf":
            with pdfplumber.open(path) as pdf:
                meta = pdf.metadata or {}
                txt  = "\n".join(p.extract_text() or "" for p in pdf.pages)
            return DocumentRecord(
                "local", path.name, meta.get("Title"), None, meta.get("Author"),
                _clean(txt), str(path), None,
                _pdf_date(meta.get("CreationDate")), path.stat().st_size
            )

        if ext == ".docx":
            doc  = docx.Document(path)
            core = doc.core_properties
            txt  = "\n".join(p.text for p in doc.paragraphs)
            return DocumentRecord(
                "local", path.name, core.title, None, core.author,
                _clean(txt), str(path), None,
                _py_dt(core.created), path.stat().st_size
            )

        if ext == ".pptx":
            prs   = Presentation(path)
            core  = prs.core_properties
            slides_txt = [
                shape.text for sl in prs.slides for shape in sl.shapes
                if hasattr(shape, "text")
            ]
            return DocumentRecord(
                "local", path.name, core.title, None, core.author,
                _clean("\n".join(slides_txt)), str(path), None,
                _py_dt(core.created), path.stat().st_size
            )

        if ext in {".html", ".htm"}:
            soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "lxml")
            return DocumentRecord(
                "local", path.name,
                soup.title.string.strip() if soup.title else None,
                None, _meta(soup, "author"),
                _clean(soup.get_text(" ", strip=True)), str(path), None,
                None, path.stat().st_size
            )

        if ext in {
            ".txt", ".md", ".markdown", ".rst", ".py", ".sh", ".r", ".R",
            ".yaml", ".yml", ".json", ".toml", ".ini", ".cfg", ".conf",
            ".csv", ".tsv", ".sql", ".xml"
        }:
            txt = path.read_text(encoding="utf-8", errors="ignore")
            return DocumentRecord(
                "local", path.name, None, None, None,
                _clean(txt), str(path), None,
                None, path.stat().st_size
            )

    except Exception as ex:
        print(f"[warn] fallback extractor failed for {path}: {ex}")

    return None
